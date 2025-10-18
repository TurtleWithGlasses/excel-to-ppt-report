"""
PowerPoint generation service.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
from typing import Dict, List, Any, Optional
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class PPTGenerator:
    """Service for generating PowerPoint presentations."""
    
    def __init__(self):
        self.presentation = None
    
    def create_presentation(self) -> Presentation:
        """
        Create a new PowerPoint presentation.
        """
        self.presentation = Presentation()
        logger.info("New presentation created")
        return self.presentation
    
    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """
        Add a title slide to the presentation.
        """
        if not self.presentation:
            self.create_presentation()
        
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        logger.info(f"Title slide added: {title}")
    
    def add_content_slide(self, title: str, content: str) -> None:
        """
        Add a content slide with title and text.
        """
        if not self.presentation:
            self.create_presentation()
        
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        text_frame = content_shape.text_frame
        text_frame.text = content
        
        logger.info(f"Content slide added: {title}")
    
    def add_table_slide(self, title: str, df: pd.DataFrame) -> None:
        """
        Add a slide with a table from DataFrame.
        """
        if not self.presentation:
            self.create_presentation()
        
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
        )
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Add table
        rows, cols = df.shape
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        table = slide.shapes.add_table(
            rows + 1, cols, left, top, width, height
        ).table
        
        # Set column headers
        for col_idx, column_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(column_name)
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Fill table with data
        for row_idx, row_data in enumerate(df.values):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value) if pd.notna(value) else ""
        
        logger.info(f"Table slide added: {title} with {rows} rows and {cols} columns")
    
    def add_chart_slide(
        self, 
        title: str, 
        df: pd.DataFrame, 
        x_column: str, 
        y_column: str,
        chart_type: str = "bar"
    ) -> None:
        """
        Add a slide with a chart from DataFrame.
        """
        if not self.presentation:
            self.create_presentation()
        
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
        )
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Prepare chart data
        chart_data = CategoryChartData()
        chart_data.categories = df[x_column].tolist()
        chart_data.add_series('Series 1', df[y_column].tolist())
        
        # Determine chart type
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }
        
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)
        
        # Add chart
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, chart_data
        ).chart
        
        logger.info(f"Chart slide added: {title} ({chart_type})")
    
    def save_presentation(self, file_path: str) -> str:
        """
        Save the presentation to a file.
        """
        if not self.presentation:
            raise ValueError("No presentation to save")
        
        # Ensure directory exists
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        
        self.presentation.save(file_path)
        logger.info(f"Presentation saved to: {file_path}")
        return file_path
    
    def generate_from_template(
        self, 
        template_structure: Dict[str, Any],
        data_sheets: Dict[str, pd.DataFrame],
        ai_insights: Optional[Dict[str, str]] = None
    ) -> None:
        """
        Generate presentation from template structure.
        """
        self.create_presentation()
        
        # Add title slide
        template_name = template_structure.get('template_name', 'Report')
        self.add_title_slide(template_name, "Generated by DataDeck")
        
        # Process sections
        sections = template_structure.get('sections', [])
        for section in sections:
            section_name = section.get('name', 'Untitled Section')
            section_type = section.get('type', 'text_analysis')
            
            if section_type == 'text_analysis':
                content = ai_insights.get(section_name, 'No insights available') if ai_insights else 'No insights available'
                self.add_content_slide(section_name, content)
            
            elif section_type == 'data_table':
                data_source = section.get('data_source', [])
                if data_source and len(data_source) > 0:
                    sheet_name = data_source[0].split('.')[0] if '.' in data_source[0] else data_source[0]
                    if sheet_name in data_sheets:
                        df = data_sheets[sheet_name]
                        self.add_table_slide(section_name, df.head(10))  # Show first 10 rows
            
            elif section_type == 'visualization':
                chart_type = section.get('chart_type', 'bar_chart').replace('_chart', '')
                data_source = section.get('data_source', [])
                if data_source and len(data_source) > 0:
                    sheet_name = data_source[0].split('.')[0] if '.' in data_source[0] else data_source[0]
                    if sheet_name in data_sheets:
                        df = data_sheets[sheet_name]
                        if len(df.columns) >= 2:
                            self.add_chart_slide(
                                section_name, 
                                df, 
                                df.columns[0], 
                                df.columns[1],
                                chart_type
                            )
        
        logger.info(f"Presentation generated from template with {len(sections)} sections")

