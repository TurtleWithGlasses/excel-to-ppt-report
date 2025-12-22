"""
ReportForge - Main Application Window (Report Generator)
Simple 4-step workflow: Import Data → Select Template → Prepare Report → Download
"""

from PyQt6.QtWidgets import (
    QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QPushButton,
    QLabel, QLineEdit, QFileDialog, QComboBox, QProgressBar,
    QGraphicsView, QGraphicsScene, QMessageBox, QFrame
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
from PyQt6.QtGui import QFont, QColor, QPainter, QPen, QPainterPath
import os
import json
import pandas as pd
from datetime import datetime

# Import core PPTGenerator
try:
    from core import PPTGenerator
    CORE_AVAILABLE = True
except ImportError:
    CORE_AVAILABLE = False
    print("Warning: Core engine not available. Using simulation mode.")


class ReportGeneratorThread(QThread):
    """Background thread for report generation"""
    progress = pyqtSignal(int, str)  # (percentage, message)
    finished = pyqtSignal(bool, str, str)  # (success, message, output_path)

    def __init__(self, excel_path, template_path, output_path, variables):
        super().__init__()
        self.excel_path = excel_path
        self.template_path = template_path
        self.output_path = output_path
        self.variables = variables

    def run(self):
        """Generate report in background"""
        try:
            if CORE_AVAILABLE:
                # Use actual PPTGenerator
                self.progress.emit(10, "Initializing generator...")
                generator = PPTGenerator()

                self.progress.emit(20, "Loading template...")
                generator.load_template(self.template_path)

                self.progress.emit(40, "Loading data...")
                generator.load_data(self.excel_path)

                self.progress.emit(60, "Setting variables...")
                generator.set_variables(self.variables)

                self.progress.emit(80, "Generating PowerPoint...")
                output = generator.generate(self.output_path)

                self.progress.emit(100, "Complete!")
                self.finished.emit(True, "Report generated successfully!", output)
            else:
                # Simulation mode
                for i in range(1, 101, 10):
                    self.progress.emit(i, f"Generating slides... {i}% complete")
                    self.msleep(200)

                self.finished.emit(True, "Report generated (simulation mode)!", self.output_path)

        except Exception as e:
            import traceback
            error_msg = f"Error: {str(e)}\n\n{traceback.format_exc()}"
            self.finished.emit(False, error_msg, "")


class StepButton(QPushButton):
    """Custom button for workflow steps"""
    def __init__(self, step_number, title, description):
        super().__init__()
        self.step_number = step_number
        self.title = title
        self.description = description
        self.completed = False
        self.setup_ui()

    def setup_ui(self):
        """Setup button appearance"""
        self.setFixedHeight(100)
        self.setMinimumWidth(200)
        self.setText(f"{self.step_number}. {self.title}\n{self.description}")
        self.setFont(QFont("Segoe UI", 10))
        self.update_style()

    def mark_completed(self):
        """Mark step as completed"""
        self.completed = True
        self.update_style()

    def mark_active(self):
        """Mark step as active"""
        self.update_style(active=True)

    def update_style(self, active=False):
        """Update button style based on state"""
        if self.completed:
            # Green for completed
            style = """
                QPushButton {
                    background-color: #10B981;
                    color: white;
                    border: 2px solid #059669;
                    border-radius: 8px;
                    padding: 10px;
                    text-align: left;
                }
                QPushButton:hover {
                    background-color: #059669;
                }
            """
        elif active:
            # Blue for active
            style = """
                QPushButton {
                    background-color: #2563EB;
                    color: white;
                    border: 2px solid #1D4ED8;
                    border-radius: 8px;
                    padding: 10px;
                    text-align: left;
                }
                QPushButton:hover {
                    background-color: #1D4ED8;
                }
            """
        else:
            # Gray for pending
            style = """
                QPushButton {
                    background-color: #F9FAFB;
                    color: #6B7280;
                    border: 2px solid #E5E7EB;
                    border-radius: 8px;
                    padding: 10px;
                    text-align: left;
                }
                QPushButton:hover {
                    background-color: #F3F4F6;
                }
            """
        self.setStyleSheet(style)


class MainWindow(QMainWindow):
    """Main Application Window - Report Generator"""

    def __init__(self):
        super().__init__()
        self.excel_path = None
        self.template_name = None
        self.template_path = None
        self.generated_slides = []
        self.current_slide_index = 0

        # Load templates dynamically from templates/configs/
        self.template_map = self.load_templates()

        self.init_ui()

    def init_ui(self):
        """Initialize user interface"""
        self.setMinimumSize(1024, 768)

        # Create central widget
        central_widget = QWidget()
        self.setCentralWidget(central_widget)

        # Main layout
        main_layout = QVBoxLayout(central_widget)
        main_layout.setSpacing(20)
        main_layout.setContentsMargins(20, 20, 20, 20)

        # Add header with Template Builder button
        self._create_header(main_layout)

        # Add components
        self._create_progress_steps(main_layout)
        self._create_separator(main_layout)
        self._create_report_name_field(main_layout)
        self._create_slide_preview(main_layout)
        self._create_slide_controls(main_layout)

    def _create_header(self, layout):
        """Create header with app title and Template Builder button"""
        header_layout = QHBoxLayout()

        # App title
        title = QLabel("📊 ReportForge - Report Generator")
        title.setFont(QFont("Segoe UI", 16, QFont.Weight.Bold))
        title.setStyleSheet("color: #1F2937;")
        header_layout.addWidget(title)

        header_layout.addStretch()

        # Template Builder button
        template_builder_btn = QPushButton("🛠️ Create/Edit Templates")
        template_builder_btn.setFont(QFont("Segoe UI", 11, QFont.Weight.Bold))
        template_builder_btn.setFixedHeight(40)
        template_builder_btn.setStyleSheet("""
            QPushButton {
                background-color: #F59E0B;
                color: white;
                border: none;
                border-radius: 6px;
                padding: 10px 20px;
            }
            QPushButton:hover {
                background-color: #D97706;
            }
        """)
        template_builder_btn.clicked.connect(self.open_template_builder)
        header_layout.addWidget(template_builder_btn)

        layout.addLayout(header_layout)

    def _create_progress_steps(self, layout):
        """Create 4-step progress workflow"""
        steps_layout = QHBoxLayout()
        steps_layout.setSpacing(10)

        # Step 1: Import Data
        self.step1_btn = StepButton(1, "Import Data", "Users will import\nexcel files")
        self.step1_btn.clicked.connect(self.import_data)
        steps_layout.addWidget(self.step1_btn)

        # Arrow
        arrow1 = QLabel("→")
        arrow1.setFont(QFont("Segoe UI", 24, QFont.Weight.Bold))
        arrow1.setAlignment(Qt.AlignmentFlag.AlignCenter)
        steps_layout.addWidget(arrow1)

        # Step 2: Select Template
        self.step2_btn = StepButton(2, "Select Template", "Users will select\ntemplate")
        self.step2_btn.clicked.connect(self.select_template)
        steps_layout.addWidget(self.step2_btn)

        # Arrow
        arrow2 = QLabel("→")
        arrow2.setFont(QFont("Segoe UI", 24, QFont.Weight.Bold))
        arrow2.setAlignment(Qt.AlignmentFlag.AlignCenter)
        steps_layout.addWidget(arrow2)

        # Step 3: Prepare Report
        self.step3_btn = StepButton(3, "Prepare Report", "Report will be prepared\nby excel importation")
        self.step3_btn.clicked.connect(self.prepare_report)
        steps_layout.addWidget(self.step3_btn)

        # Arrow
        arrow3 = QLabel("→")
        arrow3.setFont(QFont("Segoe UI", 24, QFont.Weight.Bold))
        arrow3.setAlignment(Qt.AlignmentFlag.AlignCenter)
        steps_layout.addWidget(arrow3)

        # Step 4: Download Report
        self.step4_btn = StepButton(4, "Download Report", "Report will be downloaded\non local file")
        self.step4_btn.clicked.connect(self.download_report)
        steps_layout.addWidget(self.step4_btn)

        layout.addLayout(steps_layout)

        # Progress bar (hidden initially)
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.progress_bar.setTextVisible(True)
        layout.addWidget(self.progress_bar)

    def _create_separator(self, layout):
        """Create horizontal separator line"""
        line = QFrame()
        line.setFrameShape(QFrame.Shape.HLine)
        line.setFrameShadow(QFrame.Shadow.Sunken)
        line.setStyleSheet("background-color: #E5E7EB;")
        layout.addWidget(line)

    def _create_report_name_field(self, layout):
        """Create report name input field"""
        name_layout = QHBoxLayout()

        label = QLabel("Report name:")
        label.setFont(QFont("Segoe UI", 11))
        name_layout.addWidget(label)

        self.report_name_input = QLineEdit()
        self.report_name_input.setPlaceholderText("Enter report name...")
        self.report_name_input.setText(f"Report_{datetime.now().strftime('%Y%m%d')}")
        self.report_name_input.setFont(QFont("Segoe UI", 11))
        self.report_name_input.setStyleSheet("""
            QLineEdit {
                padding: 8px;
                border: 2px solid #E5E7EB;
                border-radius: 4px;
                background-color: white;
            }
            QLineEdit:focus {
                border: 2px solid #2563EB;
            }
        """)
        name_layout.addWidget(self.report_name_input)

        layout.addLayout(name_layout)

    def _create_slide_preview(self, layout):
        """Create slide preview area"""
        # Preview container
        preview_frame = QFrame()
        preview_frame.setFrameShape(QFrame.Shape.Box)
        preview_frame.setStyleSheet("""
            QFrame {
                background-color: white;
                border: 2px solid #E5E7EB;
                border-radius: 8px;
            }
        """)
        preview_layout = QVBoxLayout(preview_frame)

        # Slide counter
        self.slide_counter = QLabel("Slide ... of ...")
        self.slide_counter.setFont(QFont("Segoe UI", 10))
        self.slide_counter.setAlignment(Qt.AlignmentFlag.AlignCenter)
        preview_layout.addWidget(self.slide_counter)

        # Graphics view for slide preview
        self.slide_view = QGraphicsView()
        self.slide_scene = QGraphicsScene()
        self.slide_view.setScene(self.slide_scene)
        self.slide_view.setRenderHint(QPainter.RenderHint.Antialiasing)
        self.slide_view.setStyleSheet("border: none; background-color: #F9FAFB;")

        # Show placeholder message
        self.show_placeholder_message()

        preview_layout.addWidget(self.slide_view)
        layout.addWidget(preview_frame)

    def _create_slide_controls(self, layout):
        """Create slide navigation and editing controls"""
        controls_layout = QHBoxLayout()
        controls_layout.setSpacing(10)

        # Previous button
        self.prev_btn = QPushButton("◄ Previous")
        self.prev_btn.setFont(QFont("Segoe UI", 10))
        self.prev_btn.setEnabled(False)
        self.prev_btn.clicked.connect(self.previous_slide)
        self.prev_btn.setStyleSheet(self._get_button_style())
        controls_layout.addWidget(self.prev_btn)

        # Edit Slide button
        self.edit_btn = QPushButton("Edit Slide")
        self.edit_btn.setFont(QFont("Segoe UI", 10))
        self.edit_btn.setEnabled(False)
        self.edit_btn.clicked.connect(self.edit_slide)
        self.edit_btn.setStyleSheet(self._get_button_style())
        controls_layout.addWidget(self.edit_btn)

        # Delete Slide button
        self.delete_btn = QPushButton("Delete Slide")
        self.delete_btn.setFont(QFont("Segoe UI", 10))
        self.delete_btn.setEnabled(False)
        self.delete_btn.clicked.connect(self.delete_slide)
        self.delete_btn.setStyleSheet(self._get_button_style("#EF4444", "#DC2626"))
        controls_layout.addWidget(self.delete_btn)

        # Add Slide button
        self.add_btn = QPushButton("Add Slide")
        self.add_btn.setFont(QFont("Segoe UI", 10))
        self.add_btn.setEnabled(False)
        self.add_btn.clicked.connect(self.add_slide)
        self.add_btn.setStyleSheet(self._get_button_style("#10B981", "#059669"))
        controls_layout.addWidget(self.add_btn)

        # Next button
        self.next_btn = QPushButton("Next ►")
        self.next_btn.setFont(QFont("Segoe UI", 10))
        self.next_btn.setEnabled(False)
        self.next_btn.clicked.connect(self.next_slide)
        self.next_btn.setStyleSheet(self._get_button_style())
        controls_layout.addWidget(self.next_btn)

        layout.addLayout(controls_layout)

    def _get_button_style(self, bg_color="#2563EB", hover_color="#1D4ED8"):
        """Get button stylesheet"""
        return f"""
            QPushButton {{
                background-color: {bg_color};
                color: white;
                border: none;
                border-radius: 4px;
                padding: 10px 20px;
                font-weight: bold;
            }}
            QPushButton:hover {{
                background-color: {hover_color};
            }}
            QPushButton:disabled {{
                background-color: #D1D5DB;
                color: #9CA3AF;
            }}
        """

    def show_placeholder_message(self):
        """Show placeholder message in slide preview"""
        self.slide_scene.clear()
        text = self.slide_scene.addText(
            "After the report is prepared,\nthe slides will be shown here\n"
            "page by page. The user will be\nable to edit the pages too.",
            QFont("Segoe UI", 14)
        )
        text.setDefaultTextColor(QColor("#EF4444"))
        text_rect = text.boundingRect()
        text.setPos(-text_rect.width()/2, -text_rect.height()/2)

    # Step 1: Import Data
    def import_data(self):
        """Import Excel file"""
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Select Excel File",
            "",
            "Excel Files (*.xlsx *.xls);;All Files (*.*)"
        )

        if file_path:
            self.excel_path = file_path
            self.step1_btn.mark_completed()
            self.step2_btn.mark_active()
            QMessageBox.information(
                self,
                "File Imported",
                f"Successfully imported:\n{os.path.basename(file_path)}"
            )

    # Step 2: Select Template
    def select_template(self):
        """Select report template"""
        if not self.excel_path:
            QMessageBox.warning(self, "No Data", "Please import Excel file first!")
            return

        # Create template selection dialog
        from PyQt6.QtWidgets import QDialog, QVBoxLayout, QDialogButtonBox

        dialog = QDialog(self)
        dialog.setWindowTitle("Select Template")
        dialog.setMinimumWidth(400)

        layout = QVBoxLayout(dialog)

        label = QLabel("Choose a template:")
        label.setFont(QFont("Segoe UI", 11))
        layout.addWidget(label)

        template_combo = QComboBox()
        template_combo.setFont(QFont("Segoe UI", 10))

        # Load templates dynamically from templates/configs/
        template_items = []
        for name in sorted(self.template_map.keys()):
            template_items.append(name)

        # Add separator and Create New option
        if template_items:
            template_items.append("---")
        template_items.append("Create New Template...")

        template_combo.addItems(template_items)
        layout.addWidget(template_combo)

        buttons = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel
        )
        buttons.accepted.connect(dialog.accept)
        buttons.rejected.connect(dialog.reject)
        layout.addWidget(buttons)

        if dialog.exec() == QDialog.DialogCode.Accepted:
            self.template_name = template_combo.currentText()
            if self.template_name == "Create New Template...":
                # Open Template Builder
                reply = QMessageBox.question(
                    self,
                    "Template Builder",
                    "Open Template Builder to create a new template?",
                    QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
                )
                if reply == QMessageBox.StandardButton.Yes:
                    self.open_template_builder()
            elif self.template_name == "---":
                # Separator, ignore
                QMessageBox.information(self, "Invalid Selection", "Please select an actual template.")
                return
            else:
                # Map template name to path
                self.template_path = self.template_map.get(self.template_name)

                if not self.template_path or not os.path.exists(self.template_path):
                    QMessageBox.warning(
                        self,
                        "Template Not Found",
                        f"Template file not found for:\n{self.template_name}\n\nPlease select a valid template."
                    )
                    return

                self.step2_btn.mark_completed()
                self.step3_btn.mark_active()
                QMessageBox.information(
                    self,
                    "Template Selected",
                    f"Selected template:\n{self.template_name}\n\nPath: {self.template_path}"
                )

    # Step 3: Prepare Report
    def prepare_report(self):
        """Generate PowerPoint report"""
        if not self.excel_path or not self.template_path:
            QMessageBox.warning(
                self,
                "Missing Information",
                "Please complete Steps 1 and 2 first!"
            )
            return

        report_name = self.report_name_input.text()
        if not report_name:
            QMessageBox.warning(self, "No Report Name", "Please enter a report name!")
            return

        # Create output directory if it doesn't exist
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)

        # Generate output path
        output_path = os.path.join(output_dir, f"{report_name}.pptx")

        # Prepare variables for text substitution
        from datetime import datetime
        now = datetime.now()
        variables = {
            'month': now.strftime('%B'),  # Full month name
            'year': now.strftime('%Y'),
            'date': now.strftime('%Y-%m-%d'),
            'report_name': report_name
        }

        # Show progress bar
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)

        # Start generation thread
        self.generator_thread = ReportGeneratorThread(
            self.excel_path,
            self.template_path,
            output_path,
            variables
        )
        self.generator_thread.progress.connect(self.update_progress)
        self.generator_thread.finished.connect(self.generation_finished)
        self.generator_thread.start()

    def update_progress(self, percentage, message):
        """Update progress bar"""
        self.progress_bar.setValue(percentage)
        self.progress_bar.setFormat(message)

    def generation_finished(self, success, message, output_path):
        """Handle report generation completion"""
        self.progress_bar.setVisible(False)

        if success:
            # Store the output path
            self.output_path = output_path

            # Load template data for preview rendering
            try:
                with open(self.template_path, 'r', encoding='utf-8') as f:
                    self.template_data = json.load(f)
            except Exception as e:
                print(f"Warning: Could not load template data: {e}")
                self.template_data = {}

            # Load Excel data for preview rendering
            try:
                self.excel_data = pd.read_excel(self.excel_path)
                # Normalize column names - strip whitespace
                self.excel_data.columns = self.excel_data.columns.str.strip()
            except Exception as e:
                print(f"Warning: Could not load Excel data: {e}")
                self.excel_data = None

            # Get actual slide count from generated PowerPoint
            try:
                from pptx import Presentation
                prs = Presentation(output_path)
                slide_count = len(prs.slides)
                self.generated_slides = [f"Slide {i}" for i in range(1, slide_count + 1)]
            except Exception as e:
                print(f"Warning: Could not read slide count from PPTX: {e}")
                # Fallback to reading template to estimate
                slide_count = len(self.template_data.get('slides', []))
                self.generated_slides = [f"Slide {i}" for i in range(1, slide_count + 1)]

            self.current_slide_index = 0

            # Update UI
            self.step3_btn.mark_completed()
            self.step4_btn.mark_active()
            self.show_slide(0)
            self.enable_slide_controls(True)

            # Show success message with file location
            QMessageBox.information(
                self,
                "Success",
                f"{message}\n\nFile saved to:\n{output_path}"
            )
        else:
            QMessageBox.critical(self, "Error", message)

    def show_slide(self, index):
        """Display slide at given index"""
        if 0 <= index < len(self.generated_slides):
            self.current_slide_index = index
            self.slide_counter.setText(
                f"Slide {index + 1} of {len(self.generated_slides)}"
            )

            # Render actual slide content
            self.slide_scene.clear()

            # Check if we have template data and render the slide
            if hasattr(self, 'template_data') and self.template_data:
                slides = self.template_data.get('slides', [])
                if index < len(slides):
                    self._render_slide_preview(slides[index], index)
                else:
                    self._render_placeholder(index)
            else:
                self._render_placeholder(index)

            # Update navigation buttons
            self.prev_btn.setEnabled(index > 0)
            self.next_btn.setEnabled(index < len(self.generated_slides) - 1)

    def enable_slide_controls(self, enabled):
        """Enable/disable slide control buttons"""
        self.edit_btn.setEnabled(enabled)
        self.delete_btn.setEnabled(enabled)
        self.add_btn.setEnabled(enabled)

    def previous_slide(self):
        """Navigate to previous slide"""
        self.show_slide(self.current_slide_index - 1)

    def next_slide(self):
        """Navigate to next slide"""
        self.show_slide(self.current_slide_index + 1)

    def edit_slide(self):
        """Edit current slide"""
        QMessageBox.information(
            self,
            "Edit Slide",
            f"Editing slide {self.current_slide_index + 1}\n\n"
            "Slide editing functionality will be implemented here."
        )

    def delete_slide(self):
        """Delete current slide"""
        reply = QMessageBox.question(
            self,
            "Delete Slide",
            f"Are you sure you want to delete slide {self.current_slide_index + 1}?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )

        if reply == QMessageBox.StandardButton.Yes:
            del self.generated_slides[self.current_slide_index]
            if self.current_slide_index >= len(self.generated_slides):
                self.current_slide_index = len(self.generated_slides) - 1
            if self.generated_slides:
                self.show_slide(self.current_slide_index)
            else:
                self.show_placeholder_message()
                self.enable_slide_controls(False)

    def add_slide(self):
        """Add new slide after current"""
        QMessageBox.information(
            self,
            "Add Slide",
            "Add new slide functionality will be implemented here.\n\n"
            "User can choose from blank slide, table slide, chart slide, etc."
        )

    # Step 4: Download Report
    def download_report(self):
        """Download generated PowerPoint report"""
        if not self.generated_slides:
            QMessageBox.warning(
                self,
                "No Report",
                "Please generate report first (Step 3)!"
            )
            return

        report_name = self.report_name_input.text()
        file_path, _ = QFileDialog.getSaveFileName(
            self,
            "Save PowerPoint Report",
            f"{report_name}.pptx",
            "PowerPoint Files (*.pptx)"
        )

        if file_path:
            # TODO: Save actual PowerPoint file
            QMessageBox.information(
                self,
                "Download Complete",
                f"Report saved successfully to:\n{file_path}\n\n"
                f"Total slides: {len(self.generated_slides)}"
            )
            self.step4_btn.mark_completed()

    # Template Builder Integration
    def open_template_builder(self):
        """Open Template Builder window"""
        from gui.template_builder import TemplateBuilder

        # Create and show Template Builder window in full-screen
        self.template_builder_window = TemplateBuilder()
        self.template_builder_window.setWindowTitle("ReportForge - Template Builder")
        self.template_builder_window.showMaximized()

        # Optional: Connect signal to refresh templates when builder closes
        self.template_builder_window.destroyed.connect(self.refresh_templates)

    def load_templates(self):
        """
        Load all templates from templates/configs/ directory.
        Returns a dictionary mapping display names to file paths.
        """
        import os
        import json

        template_map = {}
        templates_dir = os.path.join(os.getcwd(), "templates", "configs")

        # Create directory if it doesn't exist
        os.makedirs(templates_dir, exist_ok=True)

        # Scan for JSON files
        try:
            for filename in os.listdir(templates_dir):
                if filename.endswith('.json'):
                    file_path = os.path.join(templates_dir, filename)

                    # Read template to get display name
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            template_data = json.load(f)

                        # Get name from metadata (PPTGenerator format) or top-level (Template Builder format)
                        if 'metadata' in template_data:
                            display_name = template_data['metadata'].get('name', filename[:-5])
                        else:
                            display_name = template_data.get('name', filename[:-5])

                        # Use relative path
                        relative_path = os.path.join("templates", "configs", filename)
                        template_map[display_name] = relative_path

                    except Exception as e:
                        print(f"Error loading template {filename}: {e}")
                        # Use filename as fallback
                        template_map[filename[:-5]] = os.path.join("templates", "configs", filename)

        except Exception as e:
            print(f"Error scanning templates directory: {e}")

        # If no templates found, return empty dict (user will need to create templates)
        return template_map

    def refresh_templates(self):
        """Refresh template list after Template Builder closes"""
        # Reload templates from directory
        self.template_map = self.load_templates()

        # Update the dropdown in Step 1
        if hasattr(self, 'template_combo'):
            current_selection = self.template_combo.currentText()

            # Clear and repopulate
            self.template_combo.clear()

            # Group templates by industry if possible
            template_items = []
            for name in sorted(self.template_map.keys()):
                template_items.append(name)

            self.template_combo.addItems(template_items)

            # Try to restore previous selection
            index = self.template_combo.findText(current_selection)
            if index >= 0:
                self.template_combo.setCurrentIndex(index)

    # ==================== Slide Preview Rendering Methods ====================

    def _render_placeholder(self, index):
        """Render placeholder when template data not available"""
        text = self.slide_scene.addText(
            f"Slide {index + 1}\n\n[Preview of slide content will appear here]",
            QFont("Segoe UI", 16)
        )
        text_rect = text.boundingRect()
        text.setPos(-text_rect.width()/2, -text_rect.height()/2)

    def _render_slide_preview(self, slide_config, index):
        """Render slide preview based on template configuration and data"""
        INCH_TO_PIXEL = 72
        slide_type = slide_config.get('type', 'Content Slide')

        # Draw slide background (10 x 5.625 inches for 16:9)
        slide_width = 10.0 * INCH_TO_PIXEL
        slide_height = 5.625 * INCH_TO_PIXEL

        # White background
        self.slide_scene.addRect(
            0, 0, slide_width, slide_height,
            QColor("#FFFFFF"), QColor("#FFFFFF")
        )

        # Render based on slide type
        if slide_type == 'Title Slide':
            self._render_title_slide_preview(slide_config)
        elif slide_type == 'Table Slide':
            self._render_table_slide_preview(slide_config)
        elif slide_type == 'Chart Slide':
            self._render_chart_slide_preview(slide_config)
        else:
            # Generic content slide
            self._render_content_slide_preview(slide_config)

    def _render_title_slide_preview(self, slide_config):
        """Render title slide preview"""
        INCH_TO_PIXEL = 72

        # Get slide settings
        slide_settings = slide_config.get('slide_settings', {})
        title_slide_settings = slide_settings.get('title_slide', {})

        # Get title and subtitle
        title_text = title_slide_settings.get('title', 'Report Title')
        subtitle_text = title_slide_settings.get('subtitle', '')

        # Render title
        title_font = QFont(title_slide_settings.get('title_font', 'Calibri'))
        title_font.setPointSize(title_slide_settings.get('title_size', 44))
        title_font.setBold(True)
        title_item = self.slide_scene.addText(title_text, title_font)
        title_item.setDefaultTextColor(QColor(title_slide_settings.get('title_color', '#1F2937')))
        title_rect = title_item.boundingRect()
        title_item.setPos(
            5.0 * INCH_TO_PIXEL - title_rect.width() / 2,
            2.0 * INCH_TO_PIXEL
        )

        # Render subtitle if exists
        if subtitle_text:
            subtitle_font = QFont(title_slide_settings.get('subtitle_font', 'Calibri'))
            subtitle_font.setPointSize(title_slide_settings.get('subtitle_size', 20))
            subtitle_item = self.slide_scene.addText(subtitle_text, subtitle_font)
            subtitle_item.setDefaultTextColor(QColor(title_slide_settings.get('subtitle_color', '#6B7280')))
            subtitle_rect = subtitle_item.boundingRect()
            subtitle_item.setPos(
                5.0 * INCH_TO_PIXEL - subtitle_rect.width() / 2,
                3.0 * INCH_TO_PIXEL
            )

    def _render_table_slide_preview(self, slide_config):
        """Render table slide preview with actual data"""
        INCH_TO_PIXEL = 72

        # Get slide settings
        slide_settings = slide_config.get('slide_settings', {})
        table_settings = slide_settings.get('table', {})

        # Render title if exists
        title = table_settings.get('title', '')
        if title:
            title_font = QFont("Calibri", 18)
            title_font.setBold(True)
            title_item = self.slide_scene.addText(title, title_font)
            title_item.setDefaultTextColor(QColor("#1F2937"))
            title_item.setPos(0.5 * INCH_TO_PIXEL, 0.3 * INCH_TO_PIXEL)

        # Get table data
        table_data = self._get_table_data(table_settings)

        if table_data is not None and not table_data.empty:
            # Render table with data
            self._draw_data_table(table_data, table_settings)
        else:
            # Render placeholder table
            self._draw_placeholder_table(table_settings)

    def _render_chart_slide_preview(self, slide_config):
        """Render chart slide preview with actual data"""
        INCH_TO_PIXEL = 72

        # Get slide settings
        slide_settings = slide_config.get('slide_settings', {})
        chart_settings = slide_settings.get('chart', {})

        # Render title if exists
        title = chart_settings.get('title', '')
        if title:
            title_font = QFont("Calibri", 18)
            title_font.setBold(True)
            title_item = self.slide_scene.addText(title, title_font)
            title_item.setDefaultTextColor(QColor("#1F2937"))
            title_item.setPos(0.5 * INCH_TO_PIXEL, 0.3 * INCH_TO_PIXEL)

        # Get chart type and data
        chart_type = chart_settings.get('chart_type', 'column')
        chart_data = self._get_chart_data(chart_settings)

        # Chart area
        chart_y = 1.0 * INCH_TO_PIXEL
        chart_x = 0.5 * INCH_TO_PIXEL
        chart_width = 9.0 * INCH_TO_PIXEL
        chart_height = 3.5 * INCH_TO_PIXEL

        # Chart background
        self.slide_scene.addRect(
            chart_x, chart_y,
            chart_width, chart_height,
            QColor("#FFFFFF"), QColor("#E5E7EB")
        )

        # Get colors
        chart_style = chart_settings.get('style', {})
        colors = chart_style.get('colors', ['#2563EB', '#10B981', '#F59E0B', '#EF4444', '#8B5CF6'])

        # Render chart based on type
        if chart_data is not None and not chart_data.empty:
            self._draw_data_chart(chart_data, chart_type, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)
        else:
            # Render placeholder chart
            self._draw_placeholder_chart(chart_type, chart_x, chart_y, chart_width, chart_height, colors)

    def _render_content_slide_preview(self, slide_config):
        """Render generic content slide preview"""
        INCH_TO_PIXEL = 72

        slide_name = slide_config.get('name', 'Content Slide')

        # Simple placeholder for content slides
        text = self.slide_scene.addText(
            f"{slide_name}\n\n[Content slide preview]",
            QFont("Calibri", 16)
        )
        text.setDefaultTextColor(QColor("#1F2937"))
        text_rect = text.boundingRect()
        text.setPos(
            5.0 * INCH_TO_PIXEL - text_rect.width() / 2,
            2.8 * INCH_TO_PIXEL - text_rect.height() / 2
        )

    def _get_table_data(self, table_settings):
        """Get table data from Excel based on table settings"""
        if not hasattr(self, 'excel_data') or self.excel_data is None:
            return None

        try:
            # Get column mapping
            columns = table_settings.get('columns', [])
            if not columns:
                # Use all columns if not specified
                return self.excel_data.head(10)

            # Map columns - only select columns that exist in Excel
            selected_columns = []
            for col_config in columns:
                if isinstance(col_config, dict):
                    source_col = col_config.get('source_column')
                    if source_col and source_col in self.excel_data.columns:
                        selected_columns.append(source_col)
                elif isinstance(col_config, str) and col_config in self.excel_data.columns:
                    selected_columns.append(col_config)

            if not selected_columns:
                print("No valid columns found in table configuration")
                return self.excel_data.head(10)

            # Check if we need to group data (look for "Firma" or similar grouping column)
            group_by = table_settings.get('group_by')

            # If no explicit group_by, but "Firma" is in columns, group by Firma
            if not group_by and 'Firma' in selected_columns:
                group_by = 'Firma'

            # Verify group_by is both in Excel data AND in selected columns
            if group_by and group_by in self.excel_data.columns and group_by in selected_columns:
                # Group and aggregate data
                df = self.excel_data[selected_columns].copy()

                # Identify numeric columns for aggregation
                numeric_cols = df.select_dtypes(include=['number']).columns.tolist()

                # Remove group_by column from numeric cols if present
                if group_by in numeric_cols:
                    numeric_cols.remove(group_by)

                # Group by the specified column and sum numeric columns
                if numeric_cols:
                    agg_dict = {col: 'sum' for col in numeric_cols}
                    # Keep non-numeric columns by taking first value
                    for col in selected_columns:
                        if col != group_by and col not in numeric_cols:
                            agg_dict[col] = 'first'

                    df = df.groupby(group_by, as_index=False).agg(agg_dict)
                else:
                    # No numeric columns, just get unique groups
                    df = df.drop_duplicates(subset=[group_by])
            else:
                # No grouping, just select the columns
                df = self.excel_data[selected_columns].copy()

            # Apply sorting if specified
            sort_by = table_settings.get('sort_by')
            if sort_by:
                # Check if sort_by column exists, if not try to find a match
                if sort_by in df.columns:
                    ascending = table_settings.get('ascending', True)
                    df = df.sort_values(by=sort_by, ascending=ascending)
                else:
                    # Try common mappings for "Toplam" -> count or first numeric column
                    print(f"Sort column '{sort_by}' not found. Available columns: {df.columns.tolist()}")
                    # Try to sort by the first numeric column
                    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
                    if numeric_cols:
                        sort_col = numeric_cols[0]
                        print(f"Sorting by first numeric column: {sort_col}")
                        ascending = table_settings.get('ascending', True)
                        df = df.sort_values(by=sort_col, ascending=ascending)

            # Apply top_n if specified
            top_n = table_settings.get('top_n')
            if top_n and top_n > 0:
                df = df.head(top_n)
            else:
                df = df.head(10)  # Default to 10 rows

            print(f"Table data extracted: {len(df)} rows (grouped: {group_by is not None})")
            return df

        except Exception as e:
            print(f"Error getting table data: {e}")
            import traceback
            traceback.print_exc()
            return None

    def _get_chart_data(self, chart_settings):
        """Get chart data from Excel based on chart settings"""
        if not hasattr(self, 'excel_data') or self.excel_data is None:
            print("No Excel data available")
            return None

        try:
            x_column = chart_settings.get('x_column')
            y_column = chart_settings.get('y_column')

            if not x_column or not y_column:
                print(f"Missing column config: x_column={x_column}, y_column={y_column}")
                return None

            # Handle duplicate columns in DataFrame
            df_source = self.excel_data
            if df_source.columns.duplicated().any():
                print("Warning: DataFrame has duplicate column names. Removing duplicates.")
                df_source = df_source.loc[:, ~df_source.columns.duplicated()]

            if x_column not in df_source.columns:
                print(f"X column '{x_column}' not found in Excel. Available columns: {df_source.columns.tolist()}")
                return None

            # Define computed column names that need to be calculated
            computed_column_names = ['Toplam', 'Pozitif', 'Negatif', 'Nötr', 'YÜKSEK', 'ORTA', 'DÜŞÜK',
                                     'Basın', 'Radyo', 'Televizyon', 'İnternet', 'Ulusal', 'Yerel']

            # Check if y_column is a computed column
            is_computed_column = y_column in computed_column_names

            if not is_computed_column and y_column not in df_source.columns:
                print(f"Y column '{y_column}' not found in Excel. Available columns: {df_source.columns.tolist()}")
                return None

            # Handle computed columns differently
            if is_computed_column:
                # Calculate the computed column based on grouping by x_column
                df = self._calculate_computed_column_for_chart(df_source, x_column, y_column)
                if df is None or df.empty:
                    print(f"Could not calculate computed column '{y_column}'")
                    return None
            else:
                # Regular column from Excel
                # Check if y_column is numeric
                if not pd.api.types.is_numeric_dtype(df_source[y_column]):
                    print(f"Warning: Y column '{y_column}' is not numeric (dtype: {df_source[y_column].dtype})")
                    # Try to convert to numeric
                    try:
                        df_temp = df_source[[x_column, y_column]].copy()
                        df_temp[y_column] = pd.to_numeric(df_temp[y_column], errors='coerce')
                        df_temp = df_temp.dropna(subset=[y_column])
                        if df_temp.empty:
                            print(f"Could not convert '{y_column}' to numeric - all values became NaN")
                            return None
                        df = df_temp
                    except Exception as conv_err:
                        print(f"Could not convert '{y_column}' to numeric: {conv_err}")
                        return None
                else:
                    df = df_source[[x_column, y_column]].copy()

                # Apply calculation if specified
                calculation = chart_settings.get('calculation', 'sum')
                try:
                    if calculation == 'sum':
                        df = df.groupby(x_column, as_index=False)[y_column].sum()
                    elif calculation == 'mean':
                        df = df.groupby(x_column, as_index=False)[y_column].mean()
                    elif calculation == 'count':
                        df = df.groupby(x_column, as_index=False)[y_column].count()
                except Exception as groupby_err:
                    print(f"Error during groupby aggregation: {groupby_err}")
                    return None

            # Apply sorting
            sort_by = chart_settings.get('sort_by')
            if sort_by:
                ascending = chart_settings.get('ascending', True)
                if sort_by == 'x':
                    df = df.sort_values(by=x_column, ascending=ascending)
                elif sort_by == 'y':
                    df = df.sort_values(by=y_column, ascending=ascending)
                elif sort_by in df.columns:
                    df = df.sort_values(by=sort_by, ascending=ascending)
                else:
                    # Default to sorting by Y values
                    df = df.sort_values(by=y_column, ascending=ascending)

            # Apply top_n
            top_n = chart_settings.get('top_n')
            if top_n and top_n > 0:
                df = df.head(top_n)
            else:
                df = df.head(10)

            print(f"Chart data extracted successfully: {len(df)} rows")
            return df

        except Exception as e:
            print(f"Error getting chart data: {e}")
            import traceback
            traceback.print_exc()
            return None

    def _calculate_computed_column_for_chart(self, df_source, x_column, computed_column_name):
        """
        Calculate a computed column for chart data.

        Args:
            df_source: Source DataFrame from Excel
            x_column: Column to group by (X-axis)
            computed_column_name: Name of the computed column (e.g., 'Toplam', 'Basın')

        Returns:
            DataFrame with x_column and computed_column_name columns
        """
        try:
            # Handle duplicate columns
            if df_source.columns.duplicated().any():
                df_source = df_source.loc[:, ~df_source.columns.duplicated()]

            grouped = df_source.groupby(x_column)

            # Define computed column mappings
            computed_mappings = {
                'Toplam': ('__count__', None),
                'Pozitif': ('Algı', ['Pozitif', 'POZİTİF', 'POZITIF', 'pozitif']),
                'Negatif': ('Algı', ['Negatif', 'NEGATİF', 'NEGATIF', 'negatif']),
                'Nötr': ('Algı', ['Nötr', 'NÖTR', 'Notr', 'NOTR', 'nötr']),
                'YÜKSEK': ('Algı', ['YÜKSEK', 'Yüksek', 'yüksek']),
                'ORTA': ('Algı', ['ORTA', 'Orta', 'orta']),
                'DÜŞÜK': ('Algı', ['DÜŞÜK', 'Düşük', 'düşük']),
                'Basın': (['Mecra', 'Medya Tür'], ['Basın', 'BASIN', 'basın', 'Gazete', 'GAZETE']),
                'Radyo': (['Mecra', 'Medya Tür'], ['Radyo', 'RADYO', 'radyo']),
                'Televizyon': (['Mecra', 'Medya Tür'], ['Televizyon', 'TV', 'TELEVİZYON', 'televizyon']),
                'İnternet': (['Mecra', 'Medya Tür'], ['İnternet', 'İNTERNET', 'Internet', 'INTERNET', 'Online', 'ONLINE']),
                'Ulusal': ('Medya Kapsam', ['Ulusal', 'ULUSAL', 'ulusal']),
                'Yerel': ('Medya Kapsam', ['Yerel', 'YEREL', 'yerel']),
            }

            if computed_column_name not in computed_mappings:
                print(f"Unknown computed column: {computed_column_name}")
                return None

            source_col_spec, values = computed_mappings[computed_column_name]

            if source_col_spec == '__count__':
                # Total count per group
                df = grouped.size().reset_index(name=computed_column_name)
            else:
                # Find the source column (handle list of possible column names)
                if isinstance(source_col_spec, list):
                    source_col = None
                    for possible_col in source_col_spec:
                        if possible_col in df_source.columns:
                            source_col = possible_col
                            break
                else:
                    source_col = source_col_spec if source_col_spec in df_source.columns else None

                if source_col is None:
                    print(f"Source column for '{computed_column_name}' not found in data")
                    return None

                # Count specific values in the source column
                filtered = df_source[df_source[source_col].isin(values)]
                if len(filtered) > 0:
                    df = filtered.groupby(x_column).size().reset_index(name=computed_column_name)
                else:
                    # No matching values, return zeros
                    df = pd.DataFrame({
                        x_column: df_source[x_column].unique(),
                        computed_column_name: 0
                    })

            return df

        except Exception as e:
            print(f"Error calculating computed column '{computed_column_name}': {e}")
            import traceback
            traceback.print_exc()
            return None

    def _draw_data_table(self, df, table_settings):
        """Draw table with actual data"""
        INCH_TO_PIXEL = 72

        table_y = 1.0 * INCH_TO_PIXEL
        table_x = 0.5 * INCH_TO_PIXEL
        table_width = 9.0 * INCH_TO_PIXEL

        # Calculate column dimensions
        num_cols = len(df.columns)
        col_width = table_width / num_cols
        row_height = 30

        # Get style settings
        table_style = table_settings.get('style', {})
        header_color = QColor(table_style.get('header_color', '#1F2937'))
        header_text_color = QColor(table_style.get('header_text_color', '#FFFFFF'))
        row_color_1 = QColor(table_style.get('row_color_1', '#FFFFFF'))
        row_color_2 = QColor(table_style.get('row_color_2', '#F9FAFB'))
        text_color = QColor(table_style.get('text_color', '#1F2937'))
        border_color = QColor(table_style.get('border_color', '#E5E7EB'))

        # Draw header
        for col_idx, col_name in enumerate(df.columns):
            x = table_x + col_idx * col_width

            # Header cell background
            self.slide_scene.addRect(
                x, table_y, col_width, row_height,
                border_color, header_color
            )

            # Header text
            header_font = QFont("Calibri", 11)
            header_font.setBold(True)
            header_text = self.slide_scene.addText(str(col_name), header_font)
            header_text.setDefaultTextColor(header_text_color)
            header_rect = header_text.boundingRect()
            header_text.setPos(x + 5, table_y + (row_height - header_rect.height()) / 2)

        # Draw data rows (limit to 8 rows for preview)
        max_rows = min(len(df), 8)
        for row_idx in range(max_rows):
            y = table_y + (row_idx + 1) * row_height
            row_color = row_color_1 if row_idx % 2 == 0 else row_color_2

            for col_idx, col_name in enumerate(df.columns):
                x = table_x + col_idx * col_width

                # Cell background
                self.slide_scene.addRect(
                    x, y, col_width, row_height,
                    border_color, row_color
                )

                # Cell text
                value = df.iloc[row_idx][col_name]
                cell_font = QFont("Calibri", 10)
                cell_text = self.slide_scene.addText(str(value), cell_font)
                cell_text.setDefaultTextColor(text_color)
                cell_rect = cell_text.boundingRect()

                # Truncate if too long
                if cell_rect.width() > col_width - 10:
                    text_str = str(value)
                    if len(text_str) > 15:
                        text_str = text_str[:12] + "..."
                        cell_text.setPlainText(text_str)

                cell_text.setPos(x + 5, y + (row_height - cell_rect.height()) / 2)

    def _draw_placeholder_table(self, table_settings):
        """Draw placeholder table when no data available"""
        INCH_TO_PIXEL = 72

        table_y = 1.0 * INCH_TO_PIXEL
        table_x = 0.5 * INCH_TO_PIXEL
        table_width = 9.0 * INCH_TO_PIXEL
        table_height = 3.0 * INCH_TO_PIXEL

        # Draw table outline
        self.slide_scene.addRect(
            table_x, table_y,
            table_width, table_height,
            QColor("#E5E7EB"), QColor("#F9FAFB")
        )

        # Placeholder text
        text = self.slide_scene.addText("[Table data will appear here]", QFont("Calibri", 14))
        text.setDefaultTextColor(QColor("#9CA3AF"))
        text_rect = text.boundingRect()
        text.setPos(
            table_x + (table_width - text_rect.width()) / 2,
            table_y + (table_height - text_rect.height()) / 2
        )

    def _draw_data_chart(self, df, chart_type, chart_x, chart_y, chart_width, chart_height, colors, chart_settings):
        """Draw chart with actual data"""
        if chart_type == 'column':
            self._draw_column_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)
        elif chart_type == 'bar':
            self._draw_bar_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)
        elif chart_type == 'pie':
            self._draw_pie_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)
        elif chart_type == 'line':
            self._draw_line_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)
        elif chart_type == 'stacked_column':
            self._draw_stacked_column_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)
        elif chart_type == 'stacked_bar':
            self._draw_stacked_bar_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)
        else:
            self._draw_placeholder_chart(chart_type, chart_x, chart_y, chart_width, chart_height, colors)

    def _draw_column_chart_with_data(self, df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings):
        """Draw column chart with real data"""
        num_bars = min(len(df), 10)
        if num_bars == 0:
            return

        bar_width = chart_width / (num_bars + 1)

        # Get max value for scaling
        y_col = df.columns[1]
        max_val = df[y_col].max()
        if max_val == 0:
            max_val = 1

        for i in range(num_bars):
            bar_x = chart_x + (i + 0.5) * bar_width
            value = df.iloc[i][y_col]
            bar_height = (value / max_val) * chart_height * 0.8
            bar_y = chart_y + chart_height - bar_height

            color_index = i % len(colors)
            bar_color = QColor(colors[color_index])

            self.slide_scene.addRect(
                bar_x - bar_width * 0.3, bar_y,
                bar_width * 0.6, bar_height,
                bar_color, bar_color
            )

            # Value label
            value_font = QFont("Calibri", 8)
            value_text = self.slide_scene.addText(f"{value:,.0f}" if value >= 1 else f"{value:.2f}", value_font)
            value_text.setDefaultTextColor(QColor("#1F2937"))
            value_rect = value_text.boundingRect()
            value_text.setPos(bar_x - value_rect.width() / 2, bar_y - value_rect.height() - 2)

            # Category label
            x_col = df.columns[0]
            category = str(df.iloc[i][x_col])
            if len(category) > 10:
                category = category[:8] + ".."
            category_font = QFont("Calibri", 7)
            category_text = self.slide_scene.addText(category, category_font)
            category_text.setDefaultTextColor(QColor("#6B7280"))
            category_rect = category_text.boundingRect()
            category_text.setPos(bar_x - category_rect.width() / 2, chart_y + chart_height + 5)

    def _draw_bar_chart_with_data(self, df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings):
        """Draw horizontal bar chart with real data"""
        num_bars = min(len(df), 8)
        if num_bars == 0:
            return

        bar_height = chart_height / (num_bars + 1)

        # Get max value for scaling
        y_col = df.columns[1]
        max_val = df[y_col].max()
        if max_val == 0:
            max_val = 1

        for i in range(num_bars):
            bar_y = chart_y + (i + 0.5) * bar_height
            value = df.iloc[i][y_col]
            bar_width = (value / max_val) * chart_width * 0.7

            color_index = i % len(colors)
            bar_color = QColor(colors[color_index])

            self.slide_scene.addRect(
                chart_x + 100, bar_y - bar_height * 0.3,
                bar_width, bar_height * 0.6,
                bar_color, bar_color
            )

            # Value label
            value_font = QFont("Calibri", 8)
            value_text = self.slide_scene.addText(f"{value:,.0f}" if value >= 1 else f"{value:.2f}", value_font)
            value_text.setDefaultTextColor(QColor("#1F2937"))
            value_rect = value_text.boundingRect()
            value_text.setPos(chart_x + 100 + bar_width + 5, bar_y - value_rect.height() / 2)

            # Category label
            x_col = df.columns[0]
            category = str(df.iloc[i][x_col])
            if len(category) > 12:
                category = category[:10] + ".."
            category_font = QFont("Calibri", 8)
            category_text = self.slide_scene.addText(category, category_font)
            category_text.setDefaultTextColor(QColor("#6B7280"))
            category_rect = category_text.boundingRect()
            category_text.setPos(chart_x + 10, bar_y - category_rect.height() / 2)

    def _draw_pie_chart_with_data(self, df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings):
        """Draw pie chart with real data"""
        num_slices = min(len(df), 8)
        if num_slices == 0:
            return

        center_x = chart_x + chart_width / 2
        center_y = chart_y + chart_height / 2
        radius = min(chart_width, chart_height) * 0.35

        # Get values
        y_col = df.columns[1]
        total = df[y_col].sum()
        if total == 0:
            return

        # Draw slices
        start_angle = 0
        for i in range(num_slices):
            value = df.iloc[i][y_col]
            angle = (value / total) * 360

            color_index = i % len(colors)
            slice_color = QColor(colors[color_index])

            from PyQt6.QtCore import QRectF
            rect = QRectF(center_x - radius, center_y - radius, radius * 2, radius * 2)
            path = QPainterPath()
            path.moveTo(center_x, center_y)
            path.arcTo(rect, start_angle, angle)
            path.closeSubpath()

            self.slide_scene.addPath(path, slice_color, slice_color)

            start_angle += angle

    def _draw_line_chart_with_data(self, df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings):
        """Draw line chart with real data"""
        num_points = min(len(df), 10)
        if num_points == 0:
            return

        point_spacing = chart_width / (num_points + 1)

        # Get max value for scaling
        y_col = df.columns[1]
        max_val = df[y_col].max()
        min_val = df[y_col].min()
        if max_val == min_val:
            max_val = min_val + 1

        pen = QPen(QColor(colors[0]))
        pen.setWidth(3)

        # Draw lines and points
        for i in range(num_points - 1):
            x1 = chart_x + (i + 1) * point_spacing
            value1 = df.iloc[i][y_col]
            y1 = chart_y + chart_height - ((value1 - min_val) / (max_val - min_val)) * chart_height * 0.8

            x2 = chart_x + (i + 2) * point_spacing
            value2 = df.iloc[i + 1][y_col]
            y2 = chart_y + chart_height - ((value2 - min_val) / (max_val - min_val)) * chart_height * 0.8

            # Draw line
            self.slide_scene.addLine(x1, y1, x2, y2, pen)

            # Draw point
            self.slide_scene.addEllipse(x1 - 4, y1 - 4, 8, 8, pen, QColor(colors[0]))

        # Draw last point
        x_last = chart_x + num_points * point_spacing
        value_last = df.iloc[num_points - 1][y_col]
        y_last = chart_y + chart_height - ((value_last - min_val) / (max_val - min_val)) * chart_height * 0.8
        self.slide_scene.addEllipse(x_last - 4, y_last - 4, 8, 8, pen, QColor(colors[0]))

    def _draw_stacked_column_chart_with_data(self, df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings):
        """Draw stacked column chart (simplified - uses same as column for now)"""
        self._draw_column_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)

    def _draw_stacked_bar_chart_with_data(self, df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings):
        """Draw stacked bar chart (simplified - uses same as bar for now)"""
        self._draw_bar_chart_with_data(df, chart_x, chart_y, chart_width, chart_height, colors, chart_settings)

    def _draw_placeholder_chart(self, chart_type, chart_x, chart_y, chart_width, chart_height, colors):
        """Draw placeholder chart when no data available"""
        text = self.slide_scene.addText("[Chart data will appear here]", QFont("Calibri", 14))
        text.setDefaultTextColor(QColor("#9CA3AF"))
        text_rect = text.boundingRect()
        text.setPos(
            chart_x + (chart_width - text_rect.width()) / 2,
            chart_y + (chart_height - text_rect.height()) / 2
        )
