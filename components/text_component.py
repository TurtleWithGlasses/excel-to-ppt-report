"""
TextComponent - Renders text elements (titles, headings, paragraphs)

Used for: Slide titles, headers, body text, labels, static text
Supports: Font styling, colors, alignment, dynamic variables
"""

from typing import Any, Dict
from pptx.slide import Slide
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
from components.base_component import BaseComponent


class TextComponent(BaseComponent):
    """
    Component for rendering text on PowerPoint slides.

    Supports static text and dynamic variables like {date}, {company}, {reportName}.
    """

    def __init__(self, config: Dict[str, Any]):
        """
        Initialize TextComponent.

        Config keys:
            - content: Text content (string or template with {variables})
            - style:
                - font_name: Font family (default: Calibri)
                - font_size: Size in points (default: 18)
                - bold: Bold text (default: False)
                - italic: Italic text (default: False)
                - color: Hex color (default: #000000)
                - alignment: left/center/right/justify (default: left)
        """
        super().__init__(config)
        self.content = config.get('content', '')
        self.variables = config.get('variables', {})

    def validate(self) -> bool:
        """Validate text component configuration."""
        super().validate()

        if not isinstance(self.content, str):
            raise ValueError("Text content must be a string")

        return True

    def render(self, slide: Slide, data: Any = None) -> None:
        """
        Render text on the slide.

        Args:
            slide: PowerPoint slide object
            data: Optional data dict for variable substitution
        """
        # Substitute variables in content
        rendered_content = self._substitute_variables(data)

        # Add text box to slide
        text_box = slide.shapes.add_textbox(
            self.x,
            self.y,
            self.width,
            self.height
        )

        # Get text frame
        text_frame = text_box.text_frame
        text_frame.clear()  # Clear default paragraph

        # Add paragraph
        paragraph = text_frame.paragraphs[0]
        paragraph.text = rendered_content

        # Apply formatting
        self._apply_formatting(paragraph)

    def _substitute_variables(self, data: Any = None) -> str:
        """
        Replace {variables} in content with actual values.

        Args:
            data: Data dict containing variable values

        Returns:
            str: Content with substituted variables
        """
        content = self.content

        # Merge provided data with default variables
        all_vars = {**self.variables}
        if data and isinstance(data, dict):
            all_vars.update(data)

        # Replace {variable} patterns
        for key, value in all_vars.items():
            placeholder = f"{{{key}}}"
            if placeholder in content:
                content = content.replace(placeholder, str(value))

        return content

    def _apply_formatting(self, paragraph) -> None:
        """
        Apply font and paragraph formatting.

        Args:
            paragraph: python-pptx paragraph object
        """
        # Set alignment
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        alignment = self.get_alignment()
        paragraph.alignment = alignment_map.get(alignment, PP_ALIGN.LEFT)

        # Apply font formatting to all runs
        for run in paragraph.runs:
            font = run.font
            font.name = self.get_font_name()
            font.size = self.get_font_size(default=18)
            font.bold = self.get_bold()
            font.italic = self.get_italic()

            # Set color
            r, g, b = self.get_color('color', '#000000')
            font.color.rgb = RGBColor(r, g, b)

    def set_content(self, content: str) -> None:
        """
        Update text content.

        Args:
            content: New text content
        """
        self.content = content

    def add_variable(self, key: str, value: Any) -> None:
        """
        Add a variable for substitution.

        Args:
            key: Variable name (used as {key} in content)
            value: Variable value
        """
        self.variables[key] = value

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        config = super().to_dict()
        config['content'] = self.content
        config['variables'] = self.variables
        return config


# Example usage configurations:
EXAMPLE_CONFIGS = {
    'title': {
        'type': 'text',
        'content': 'Monthly Media Report - {month} {year}',
        'position': {'x': 0.5, 'y': 0.5},
        'size': {'width': 9.0, 'height': 1.0},
        'variables': {'month': 'October', 'year': '2025'},
        'style': {
            'font_name': 'Calibri',
            'font_size': 32,
            'bold': True,
            'color': '#1F2937',
            'alignment': 'center'
        }
    },
    'subtitle': {
        'type': 'text',
        'content': 'Company: {company} | Report Date: {date}',
        'position': {'x': 0.5, 'y': 1.5},
        'size': {'width': 9.0, 'height': 0.5},
        'variables': {'company': 'BSH', 'date': '2025-11-30'},
        'style': {
            'font_name': 'Calibri',
            'font_size': 16,
            'color': '#6B7280',
            'alignment': 'center'
        }
    },
    'body': {
        'type': 'text',
        'content': 'This report provides comprehensive analysis of media coverage for {company}.',
        'position': {'x': 1.0, 'y': 3.0},
        'size': {'width': 8.0, 'height': 2.0},
        'style': {
            'font_name': 'Calibri',
            'font_size': 14,
            'color': '#000000',
            'alignment': 'left'
        }
    }
}
