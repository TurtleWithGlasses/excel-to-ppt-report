"""
SummaryComponent - Generates automatic insights and summaries from data

Used for: Key metrics, highlights, trends, auto-generated insights
Supports: Statistical analysis, top performers, comparisons, AI integration (future)
"""

from typing import Any, Dict, List, Optional, Tuple
from pptx.slide import Slide
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from components.base_component import BaseComponent
import pandas as pd
import numpy as np
from datetime import datetime


class SummaryComponent(BaseComponent):
    """
    Component for generating automatic summaries and insights from data.

    Analyzes data and generates key metrics, trends, and highlights automatically.
    Future: AI-powered insights using Claude/OpenAI APIs.
    """

    # Insight types
    INSIGHT_KEY_METRICS = 'key_metrics'
    INSIGHT_TRENDS = 'trends'
    INSIGHT_HIGHLIGHTS = 'highlights'
    INSIGHT_COMPARISONS = 'comparisons'
    INSIGHT_TOP_PERFORMERS = 'top_performers'

    def __init__(self, config: Dict[str, Any]):
        """
        Initialize SummaryComponent.

        Config keys:
            - data_source:
                - insight_types: List of insights to generate
                - metric_columns: Columns to analyze for metrics
                - compare_column: Column for comparisons
                - time_column: Column for trend analysis
                - max_items: Maximum number of items to show
            - style:
                - layout: 'bullets' | 'numbered' | 'callout_boxes'
                - show_icons: Show emoji/icons (default: True)
                - highlight_color: Color for important values
                - font_size: Text size
        """
        super().__init__(config)

        # Insight configuration
        self.insight_types = self.data_source.get('insight_types', [self.INSIGHT_KEY_METRICS])
        self.metric_columns = self.data_source.get('metric_columns', [])
        self.compare_column = self.data_source.get('compare_column')
        self.time_column = self.data_source.get('time_column')
        self.max_items = self.data_source.get('max_items', 5)

        # Style configuration
        self.layout = self.style.get('layout', 'bullets')
        self.show_icons = self.style.get('show_icons', True)
        self.highlight_color = self.style.get('highlight_color', '#2563EB')

    def validate(self) -> bool:
        """Validate summary component configuration."""
        super().validate()

        if not self.insight_types:
            raise ValueError("At least one insight_type is required")

        valid_types = [
            self.INSIGHT_KEY_METRICS, self.INSIGHT_TRENDS,
            self.INSIGHT_HIGHLIGHTS, self.INSIGHT_COMPARISONS,
            self.INSIGHT_TOP_PERFORMERS
        ]

        for insight_type in self.insight_types:
            if insight_type not in valid_types:
                raise ValueError(f"Invalid insight_type: {insight_type}")

        return True

    def render(self, slide: Slide, data: Any) -> None:
        """
        Render summary insights on the slide.

        Args:
            slide: PowerPoint slide object
            data: pandas DataFrame or dict with metrics
        """
        # Prepare data
        df = self._prepare_data(data)

        # Generate insights
        insights = self._generate_insights(df, data)

        if not insights:
            self._render_placeholder(slide)
            return

        # Render based on layout
        if self.layout == 'callout_boxes':
            self._render_callout_boxes(slide, insights)
        elif self.layout == 'numbered':
            self._render_numbered_list(slide, insights)
        else:  # bullets
            self._render_bullet_list(slide, insights)

    def _prepare_data(self, data: Any) -> Optional[pd.DataFrame]:
        """Convert data to DataFrame for analysis."""
        if data is None:
            return None

        if isinstance(data, pd.DataFrame):
            return data.copy()
        elif isinstance(data, list) and data:
            return pd.DataFrame(data)
        elif isinstance(data, dict):
            # Check if it's metrics dict or data dict
            first_val = next(iter(data.values()), None)
            if isinstance(first_val, (list, pd.Series)):
                return pd.DataFrame(data)
            # Treat as single-row metrics
            return pd.DataFrame([data])

        return None

    def _generate_insights(self, df: Optional[pd.DataFrame], raw_data: Any) -> List[str]:
        """
        Generate insights from data.

        Args:
            df: DataFrame with data
            raw_data: Original data (may include metadata)

        Returns:
            List of insight strings
        """
        insights = []

        for insight_type in self.insight_types:
            if insight_type == self.INSIGHT_KEY_METRICS:
                insights.extend(self._generate_key_metrics(df, raw_data))
            elif insight_type == self.INSIGHT_TRENDS:
                insights.extend(self._generate_trends(df))
            elif insight_type == self.INSIGHT_HIGHLIGHTS:
                insights.extend(self._generate_highlights(df))
            elif insight_type == self.INSIGHT_COMPARISONS:
                insights.extend(self._generate_comparisons(df))
            elif insight_type == self.INSIGHT_TOP_PERFORMERS:
                insights.extend(self._generate_top_performers(df))

        return insights[:self.max_items]

    def _generate_key_metrics(self, df: Optional[pd.DataFrame], raw_data: Any) -> List[str]:
        """Generate key metrics insights."""
        insights = []
        icon = "📊 " if self.show_icons else ""

        # If raw_data is a metrics dict, use it directly
        if isinstance(raw_data, dict) and df is None:
            for key, value in list(raw_data.items())[:3]:
                if isinstance(value, (int, float)):
                    formatted_value = f"{value:,.0f}" if isinstance(value, int) else f"{value:,.2f}"
                    insights.append(f"{icon}{key}: {formatted_value}")
            return insights

        if df is None or df.empty:
            return insights

        # Analyze numeric columns
        for col in self.metric_columns or df.select_dtypes(include=[np.number]).columns:
            if col not in df.columns:
                continue

            total = df[col].sum()
            avg = df[col].mean()
            max_val = df[col].max()

            insights.append(f"{icon}Total {col}: {total:,.0f}")
            insights.append(f"{icon}Average {col}: {avg:,.1f}")

            if len(insights) >= 3:
                break

        return insights

    def _generate_trends(self, df: Optional[pd.DataFrame]) -> List[str]:
        """Generate trend insights."""
        insights = []
        icon = "📈 " if self.show_icons else ""

        if df is None or df.empty or not self.time_column:
            return insights

        if self.time_column not in df.columns:
            return insights

        # Sort by time
        df_sorted = df.sort_values(by=self.time_column)

        # Analyze trends for numeric columns
        for col in self.metric_columns or df_sorted.select_dtypes(include=[np.number]).columns[:2]:
            if col == self.time_column or col not in df_sorted.columns:
                continue

            values = df_sorted[col].values
            if len(values) < 2:
                continue

            # Calculate trend
            first_val = values[0]
            last_val = values[-1]
            change = last_val - first_val
            pct_change = (change / first_val * 100) if first_val != 0 else 0

            trend_icon = "📈" if change > 0 else "📉" if change < 0 else "➡️"
            trend_word = "increased" if change > 0 else "decreased" if change < 0 else "remained stable"

            if self.show_icons:
                insights.append(
                    f"{trend_icon} {col} {trend_word} by {abs(pct_change):.1f}% "
                    f"({first_val:,.0f} → {last_val:,.0f})"
                )
            else:
                insights.append(
                    f"{col} {trend_word} by {abs(pct_change):.1f}% over period"
                )

        return insights

    def _generate_highlights(self, df: Optional[pd.DataFrame]) -> List[str]:
        """Generate highlight insights (max/min values)."""
        insights = []
        icon = "⭐ " if self.show_icons else ""

        if df is None or df.empty:
            return insights

        for col in self.metric_columns or df.select_dtypes(include=[np.number]).columns[:2]:
            if col not in df.columns:
                continue

            max_idx = df[col].idxmax()
            max_val = df[col].max()

            # Get identifier (company name, category, etc.)
            identifier_col = self.compare_column or df.columns[0]
            if identifier_col in df.columns:
                identifier = df.loc[max_idx, identifier_col]
                insights.append(f"{icon}Highest {col}: {identifier} ({max_val:,.0f})")
            else:
                insights.append(f"{icon}Highest {col}: {max_val:,.0f}")

        return insights

    def _generate_comparisons(self, df: Optional[pd.DataFrame]) -> List[str]:
        """Generate comparison insights."""
        insights = []
        icon = "🔄 " if self.show_icons else ""

        if df is None or df.empty or len(df) < 2:
            return insights

        if not self.compare_column or self.compare_column not in df.columns:
            return insights

        # Compare top 2 entities
        for col in self.metric_columns or df.select_dtypes(include=[np.number]).columns[:1]:
            if col not in df.columns:
                continue

            top_2 = df.nlargest(2, col)
            if len(top_2) < 2:
                continue

            first = top_2.iloc[0]
            second = top_2.iloc[1]

            first_name = first[self.compare_column]
            second_name = second[self.compare_column]
            first_val = first[col]
            second_val = second[col]

            diff = first_val - second_val
            pct_diff = (diff / second_val * 100) if second_val != 0 else 0

            insights.append(
                f"{icon}{first_name} leads {second_name} by {diff:,.0f} "
                f"({pct_diff:.1f}%) in {col}"
            )

        return insights

    def _generate_top_performers(self, df: Optional[pd.DataFrame]) -> List[str]:
        """Generate top performers insights."""
        insights = []
        icon = "🏆 " if self.show_icons else ""

        if df is None or df.empty:
            return insights

        identifier_col = self.compare_column or df.columns[0]
        if identifier_col not in df.columns:
            return insights

        for col in self.metric_columns or df.select_dtypes(include=[np.number]).columns[:1]:
            if col not in df.columns:
                continue

            top_3 = df.nlargest(3, col)
            top_names = top_3[identifier_col].tolist()

            insights.append(f"{icon}Top 3 in {col}: {', '.join(map(str, top_names))}")

        return insights

    def _render_bullet_list(self, slide: Slide, insights: List[str]) -> None:
        """Render insights as bullet list."""
        text_box = slide.shapes.add_textbox(
            self.x, self.y, self.width, self.height
        )

        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        for idx, insight in enumerate(insights):
            if idx == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()

            paragraph.text = insight
            paragraph.level = 0
            paragraph.font.name = self.get_font_name()
            paragraph.font.size = self.get_font_size(default=14)
            paragraph.font.color.rgb = RGBColor(31, 41, 55)
            paragraph.space_after = Pt(8)

    def _render_numbered_list(self, slide: Slide, insights: List[str]) -> None:
        """Render insights as numbered list."""
        text_box = slide.shapes.add_textbox(
            self.x, self.y, self.width, self.height
        )

        text_frame = text_box.text_frame
        text_frame.clear()

        for idx, insight in enumerate(insights):
            if idx == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()

            # Remove icon if present (will use number instead)
            clean_insight = insight
            if self.show_icons and any(emoji in insight for emoji in ['📊', '📈', '📉', '⭐', '🔄', '🏆', '➡️']):
                # Remove first emoji
                for emoji in ['📊 ', '📈 ', '📉 ', '⭐ ', '🔄 ', '🏆 ', '➡️ ']:
                    clean_insight = clean_insight.replace(emoji, '', 1)

            paragraph.text = f"{idx + 1}. {clean_insight}"
            paragraph.font.name = self.get_font_name()
            paragraph.font.size = self.get_font_size(default=14)
            paragraph.space_after = Pt(8)

    def _render_callout_boxes(self, slide: Slide, insights: List[str]) -> None:
        """Render insights as callout boxes."""
        from pptx.enum.shapes import MSO_SHAPE

        box_height = (self.height.inches - 0.2 * (len(insights) - 1)) / len(insights)
        box_height = Inches(max(0.5, box_height))

        for idx, insight in enumerate(insights):
            y_pos = self.y + Inches(idx * (box_height.inches + 0.2))

            # Create box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                self.x,
                y_pos,
                self.width,
                box_height
            )

            # Box styling
            r, g, b = self.get_color('highlight_color', '#EFF6FF')
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(r, g, b)

            r, g, b = self.get_color('highlight_color', '#2563EB')
            box.line.color.rgb = RGBColor(r, g, b)
            box.line.width = Pt(2)

            # Add text
            text_frame = box.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.text = insight
            paragraph.font.name = self.get_font_name()
            paragraph.font.size = self.get_font_size(default=13)
            paragraph.font.bold = True
            text_frame.vertical_anchor = 1  # Middle
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)

    def _render_placeholder(self, slide: Slide) -> None:
        """Render placeholder when no insights available."""
        text_box = slide.shapes.add_textbox(
            self.x, self.y, self.width, self.height
        )

        text_frame = text_box.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.text = "[No insights available - insufficient data]"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(156, 163, 175)
        text_frame.vertical_anchor = 1

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        config = super().to_dict()
        config['data_source'].update({
            'insight_types': self.insight_types,
            'metric_columns': self.metric_columns,
            'compare_column': self.compare_column,
            'time_column': self.time_column,
            'max_items': self.max_items
        })
        return config


# Example configurations:
EXAMPLE_CONFIGS = {
    'key_metrics': {
        'type': 'summary',
        'position': {'x': 0.5, 'y': 5.0},
        'size': {'width': 9.0, 'height': 1.5},
        'data_source': {
            'insight_types': ['key_metrics'],
            'metric_columns': ['Total', 'Positive', 'Negative'],
            'max_items': 3
        },
        'style': {
            'layout': 'bullets',
            'show_icons': True,
            'font_size': 14
        }
    },
    'trends_and_highlights': {
        'type': 'summary',
        'position': {'x': 0.5, 'y': 5.0},
        'size': {'width': 9.0, 'height': 2.0},
        'data_source': {
            'insight_types': ['trends', 'highlights'],
            'metric_columns': ['Total'],
            'compare_column': 'Company',
            'time_column': 'Month',
            'max_items': 5
        },
        'style': {
            'layout': 'callout_boxes',
            'show_icons': True,
            'highlight_color': '#EFF6FF'
        }
    }
}
