"""
TableComponent - Renders data tables from Excel or structured data

Used for: Executive summaries, data comparisons, KPI tables, metrics
Supports: Custom styling, zebra striping, headers, sorting, formatting
"""

from typing import Any, Dict, List, Optional
from pptx.slide import Slide
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches as In
from pptx.dml.color import RGBColor
from components.base_component import BaseComponent
import pandas as pd


class TableComponent(BaseComponent):
    """
    Component for rendering data tables in PowerPoint slides.

    Supports pandas DataFrames, lists of dicts, and structured data.
    """

    def __init__(self, config: Dict[str, Any], template: Optional[Dict[str, Any]] = None):
        """
        Initialize TableComponent.

        Config keys:
            - data_source:
                - type: 'dataframe' | 'list' | 'dict'
                - columns: List of column names to include
                - sort_by: Column to sort by (optional)
                - ascending: Sort order (default: True)
            - style:
                - header_row: Show header row (default: True)
                - zebra_striping: Alternate row colors (default: True)
                - border: Show table border (default: True)
                - grid_lines: Show internal grid (default: False)
                - header_color: Hex color for header bg
                - header_text_color: Hex color for header text
                - row_color_1: First row color
                - row_color_2: Second row color (for zebra)
                - text_color: Table text color
                - font_size: Table font size
        """
        super().__init__(config, template=template)

        # Table-specific config
        self.show_header = self.style.get('header_row', True)
        self.zebra_striping = self.style.get('zebra_striping', True)
        self.show_border = self.style.get('border', True)
        self.show_grid = self.style.get('grid_lines', False)

        # Column configuration
        self.columns = self.data_source.get('columns', [])
        self.column_mapping = self.data_source.get('column_mapping', {})
        self.sort_by = self.data_source.get('sort_by')
        self.ascending = self.data_source.get('ascending', True)

        # Validate after all attributes are set
        self.validate()

    def validate(self) -> bool:
        """Validate table component configuration."""
        super().validate()

        if not self.columns and not self.column_mapping:
            raise ValueError("TableComponent requires either 'columns' or 'column_mapping'")

        return True

    def render(self, slide: Slide, data: Any) -> None:
        """
        Render table on the slide.

        Args:
            slide: PowerPoint slide object
            data: pandas DataFrame, list of dicts, or dict
        """
        # Debug: Print style settings being used
        print("\n[TableComponent] ===== RENDER DEBUG =====")
        print(f"[TableComponent] Style settings received: {self.style}")
        print(f"[TableComponent] header_alignment: {self.style.get('header_alignment', 'NOT FOUND')}")
        print(f"[TableComponent] text_alignment: {self.style.get('text_alignment', 'NOT FOUND')}")
        print(f"[TableComponent] header_bold: {self.style.get('header_bold', 'NOT FOUND')}")
        print(f"[TableComponent] text_bold: {self.style.get('text_bold', 'NOT FOUND')}")

        # Convert data to DataFrame if needed
        df = self._prepare_data(data)

        if df is None or df.empty:
            # Add placeholder text if no data
            self._render_empty_table(slide)
            return

        # Create table shape
        rows = len(df) + (1 if self.show_header else 0)
        cols = len(df.columns)

        # Calculate appropriate height based on number of rows
        # For 16:9 slides (10" x 5.625"), with title at y=0.3" (height 0.6") and table at y=1.0"
        # Available space: 5.625 - 1.0 = 4.625" (need to leave small margin at bottom)
        max_height = 4.5  # Conservative to ensure all rows are visible within slide

        # Ideal row heights
        ideal_header_height = 0.35 if self.show_header else 0
        ideal_data_row_height = 0.25
        ideal_total_height = ideal_header_height + (len(df) * ideal_data_row_height)

        # If ideal height exceeds max, scale down the row heights proportionally
        if ideal_total_height > max_height:
            scale_factor = max_height / ideal_total_height
            header_height = ideal_header_height * scale_factor
            data_row_height = ideal_data_row_height * scale_factor
            final_height = max_height
            print(f"[TableComponent] Scaling down row heights by {scale_factor:.3f} to fit within {max_height}\"")
        else:
            header_height = ideal_header_height
            data_row_height = ideal_data_row_height
            final_height = ideal_total_height

        print(f"[TableComponent] Creating table with {rows} rows (data: {len(df)}, header: {self.show_header}) x {cols} cols")
        print(f"[TableComponent] Ideal height: {ideal_total_height:.2f}\", final: {final_height:.2f}\"")
        print(f"[TableComponent] Header height: {header_height:.3f}\", Data row height: {data_row_height:.3f}\"")
        print(f"[TableComponent] Table position: y={self.y}, total extent: {self.y.inches + final_height:.2f}\"")

        table_shape = slide.shapes.add_table(
            rows, cols, self.x, self.y, self.width, In(final_height)
        )
        table = table_shape.table

        # Explicitly set row heights to ensure all rows are visible
        print(f"[TableComponent] Setting explicit row heights for {rows} rows")
        for row_idx in range(rows):
            if row_idx == 0 and self.show_header:
                # Header row
                table.rows[row_idx].height = In(header_height)
                print(f"[TableComponent] Row {row_idx} (header): height = {header_height:.3f}\"")
            else:
                # Data row
                table.rows[row_idx].height = In(data_row_height)
                if row_idx < 3:  # Only print first few to avoid spam
                    print(f"[TableComponent] Row {row_idx} (data): height = {data_row_height:.3f}\"")

        # Render header
        if self.show_header:
            self._render_header(table, df.columns)

        # Render data rows
        self._render_rows(table, df)

        # Apply styling
        self._apply_table_styling(table)

        print("[TableComponent] ===== RENDER COMPLETE =====\n")

    def _prepare_data(self, data: Any) -> Optional[pd.DataFrame]:
        """
        Convert input data to pandas DataFrame.

        Args:
            data: Input data (DataFrame, list, dict, or file path)

        Returns:
            pd.DataFrame or None
        """
        if data is None:
            return None

        # Already a DataFrame
        if isinstance(data, pd.DataFrame):
            df = data.copy()

        # List of dicts
        elif isinstance(data, list):
            if not data:
                return None
            df = pd.DataFrame(data)

        # Dict (single row or column-oriented)
        elif isinstance(data, dict):
            # Check if it's column-oriented {col1: [vals], col2: [vals]}
            first_val = next(iter(data.values()))
            if isinstance(first_val, list):
                df = pd.DataFrame(data)
            else:
                # Single row dict
                df = pd.DataFrame([data])

        # File path (Excel)
        elif isinstance(data, str) and data.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(data)

        else:
            raise ValueError(f"Unsupported data type: {type(data)}")

        # Apply column filtering
        if self.columns:
            available_cols = [col for col in self.columns if col in df.columns]
            df = df[available_cols]

        # Apply column mapping (rename)
        if self.column_mapping:
            df = df.rename(columns=self.column_mapping)

        # Apply sorting
        if self.sort_by and self.sort_by in df.columns:
            df = df.sort_values(by=self.sort_by, ascending=self.ascending)

        return df

    def _render_header(self, table, columns: List[str]) -> None:
        """
        Render table header row.

        Args:
            table: python-pptx table object
            columns: List of column names
        """
        print(f"[TableComponent] Rendering header with {len(columns)} columns")

        for col_idx, col_name in enumerate(columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)

            # Header formatting
            paragraph = cell.text_frame.paragraphs[0]

            # Apply header alignment from style
            header_align = self.style.get('header_alignment', 'Center')
            print(f"[TableComponent] Header col {col_idx}: alignment = {header_align}")

            if header_align.lower() == 'left':
                paragraph.alignment = PP_ALIGN.LEFT
            elif header_align.lower() == 'right':
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.CENTER

            font = paragraph.runs[0].font
            font.name = self.get_font_name()
            font.size = self.get_font_size(default=11)

            # Apply header bold/italic from style
            header_bold = self.style.get('header_bold', True)
            header_italic = self.style.get('header_italic', False)
            print(f"[TableComponent] Header col {col_idx}: bold={header_bold}, italic={header_italic}")

            font.bold = header_bold
            font.italic = header_italic

            # Header colors
            r, g, b = self.get_color('header_color', '#2563EB')
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(r, g, b)

            r, g, b = self.get_color('header_text_color', '#FFFFFF')
            font.color.rgb = RGBColor(r, g, b)

    def _render_rows(self, table, df: pd.DataFrame) -> None:
        """
        Render data rows in table.

        Args:
            table: python-pptx table object
            df: DataFrame with data to render
        """
        row_offset = 1 if self.show_header else 0
        print(f"[TableComponent] Rendering {len(df)} data rows (row_offset={row_offset})")

        # Get text styling once (same for all cells)
        text_align = self.style.get('text_alignment', 'Left')
        text_bold = self.style.get('text_bold', False)
        text_italic = self.style.get('text_italic', False)
        print(f"[TableComponent] Text style: alignment={text_align}, bold={text_bold}, italic={text_italic}")

        for row_idx, (_, row) in enumerate(df.iterrows()):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + row_offset, col_idx)
                cell.text = self._format_value(value)

                # Cell formatting
                paragraph = cell.text_frame.paragraphs[0]

                # Apply text alignment from style
                if text_align.lower() == 'center':
                    paragraph.alignment = PP_ALIGN.CENTER
                elif text_align.lower() == 'right':
                    paragraph.alignment = PP_ALIGN.RIGHT
                else:
                    paragraph.alignment = PP_ALIGN.LEFT

                font = paragraph.runs[0].font
                font.name = self.get_font_name()
                font.size = self.get_font_size(default=10)

                # Apply text bold/italic from style
                font.bold = text_bold
                font.italic = text_italic

                # Text color
                r, g, b = self.get_color('text_color', '#000000')
                font.color.rgb = RGBColor(r, g, b)

                # Zebra striping
                if self.zebra_striping and row_idx % 2 == 1:
                    r, g, b = self.get_color('row_color_2', '#F3F4F6')
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(r, g, b)
                elif self.zebra_striping:
                    r, g, b = self.get_color('row_color_1', '#FFFFFF')
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(r, g, b)

        print(f"[TableComponent] Finished rendering {len(df)} data rows")

    def _format_value(self, value: Any) -> str:
        """
        Format cell value for display.

        Args:
            value: Cell value

        Returns:
            str: Formatted string
        """
        if pd.isna(value):
            return ""

        if isinstance(value, (int, float)):
            # Check for percentage format
            if 'number_format' in self.style:
                fmt = self.style['number_format']
                if fmt == 'percentage':
                    return f"{value:.1f}%"
                elif fmt == 'currency':
                    return f"₺{value:,.0f}"
                elif fmt == 'decimal':
                    return f"{value:.2f}"

            # Default number formatting
            if isinstance(value, int):
                return f"{value:,}"
            else:
                return f"{value:,.2f}"

        return str(value)

    def _apply_table_styling(self, table) -> None:
        """
        Apply border and grid styling to table.

        Args:
            table: python-pptx table object
        """
        # Note: python-pptx has limited border control
        # Borders are typically controlled at cell level
        pass

    def _render_empty_table(self, slide: Slide) -> None:
        """
        Render placeholder when no data is available.

        Args:
            slide: PowerPoint slide object
        """
        text_box = slide.shapes.add_textbox(
            self.x, self.y, self.width, self.height
        )
        text_frame = text_box.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.text = "[No data available for table]"
        paragraph.alignment = PP_ALIGN.CENTER

        font = paragraph.runs[0].font
        font.size = Pt(14)
        font.color.rgb = RGBColor(156, 163, 175)  # Gray

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        config = super().to_dict()
        config['data_source']['columns'] = self.columns
        config['data_source']['column_mapping'] = self.column_mapping
        config['data_source']['sort_by'] = self.sort_by
        return config


# Example usage configurations:
EXAMPLE_CONFIGS = {
    'executive_summary': {
        'type': 'table',
        'position': {'x': 0.5, 'y': 2.0},
        'size': {'width': 9.0, 'height': 4.0},
        'data_source': {
            'columns': ['Company', 'Total', 'Positive', 'Negative', 'Neutral', 'SoV'],
            'column_mapping': {
                'Kurum': 'Company',
                'Toplam': 'Total',
                'Olumlu': 'Positive',
                'Olumsuz': 'Negative',
                'Nötr': 'Neutral'
            },
            'sort_by': 'Total',
            'ascending': False
        },
        'style': {
            'header_row': True,
            'zebra_striping': True,
            'border': True,
            'grid_lines': False,
            'header_color': '#2563EB',
            'header_text_color': '#FFFFFF',
            'row_color_1': '#FFFFFF',
            'row_color_2': '#F9FAFB',
            'text_color': '#1F2937',
            'font_name': 'Calibri',
            'font_size': 10
        }
    }
}
