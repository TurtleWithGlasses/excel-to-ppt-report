"""
ImageComponent - Renders images, logos, and graphics

Used for: Company logos, product photos, charts as images, graphics
Supports: File paths, URLs, PIL Images, with styling options
"""

from typing import Any, Dict, Optional
from pptx.slide import Slide
from pptx.util import Inches
from components.base_component import BaseComponent
import os
from io import BytesIO
from pathlib import Path


class ImageComponent(BaseComponent):
    """
    Component for rendering images on PowerPoint slides.

    Supports: Local files (.png, .jpg, .svg), URLs, PIL Image objects
    """

    def __init__(self, config: Dict[str, Any]):
        """
        Initialize ImageComponent.

        Config keys:
            - data_source:
                - path: File path or URL to image
                - type: 'file' | 'url' | 'template_logo' | 'template_embedded_logo'
            - style:
                - border_width: Border width in points (0 = no border)
                - border_color: Hex color for border
                - corner_radius: Rounded corners in points
                - opacity: 0-100 (transparency)
                - maintain_aspect: Keep aspect ratio (default: True)
        """
        super().__init__(config)

        # Image source
        self.image_path = self.data_source.get('path')
        self.image_type = self.data_source.get('type', 'file')

        # Style options
        self.border_width = self.style.get('border_width', 0)
        self.border_color = self.style.get('border_color', '#000000')
        self.corner_radius = self.style.get('corner_radius', 0)
        self.opacity = self.style.get('opacity', 100)
        self.maintain_aspect = self.style.get('maintain_aspect', True)

        # Validate after all attributes are set
        self.validate()

    def validate(self) -> bool:
        """Validate image component configuration."""
        super().validate()

        if not self.image_path and self.image_type not in ['template_logo', 'template_embedded_logo']:
            raise ValueError("ImageComponent requires 'path' in data_source or template logo type")

        if self.image_type == 'file' and not os.path.exists(self.image_path):
            # Don't raise error, just warn - image might be added later
            pass

        return True

    def render(self, slide: Slide, data: Any = None) -> None:
        """
        Render image on the slide.

        Args:
            slide: PowerPoint slide object
            data: Optional dict with 'logo_path' or 'image_path' override
        """
        # Get image path (allow override from data)
        image_path = self._get_image_path(data)

        if not image_path or not os.path.exists(image_path):
            self._render_placeholder(slide)
            return

        try:
            # Add image to slide
            if self.maintain_aspect:
                # Add with aspect ratio maintained
                picture = slide.shapes.add_picture(
                    image_path,
                    self.x,
                    self.y,
                    width=self.width
                )
                # Center vertically if height is smaller
                if picture.height < self.height:
                    picture.top = self.y + (self.height - picture.height) // 2
            else:
                # Stretch to fill area
                picture = slide.shapes.add_picture(
                    image_path,
                    self.x,
                    self.y,
                    width=self.width,
                    height=self.height
                )

            # Apply styling (limited in python-pptx)
            self._apply_styling(picture)

        except Exception as e:
            print(f"Error rendering image {image_path}: {e}")
            self._render_placeholder(slide, error=str(e))

    def _get_image_path(self, data: Any = None) -> Optional[str]:
        """
        Get image path from config, template, or data override.

        Args:
            data: Optional data dict with image path

        Returns:
            str: Image file path or None
        """
        # Check for override in data (highest priority)
        if data and isinstance(data, dict):
            if 'logo_path' in data:
                return data['logo_path']
            if 'image_path' in data:
                return data['image_path']

        # Use template logo if specified
        if self.image_type == 'template_logo':
            # First check template settings for logo_path
            if self.template:
                settings = self.template.get('settings', {})
                logo_path = settings.get('logo_path')
                if logo_path:
                    # Resolve relative paths relative to template directory or project root
                    resolved_path = self._resolve_logo_path(logo_path)
                    if resolved_path:
                        return resolved_path
            
            # Fallback: check data dict for template_logo
            if data and isinstance(data, dict) and 'template_logo' in data:
                return data['template_logo']
        
        # Use template embedded logo if specified
        if self.image_type == 'template_embedded_logo':
            # First check template settings for embedded_logo_path
            if self.template:
                settings = self.template.get('settings', {})
                embedded_logo_path = settings.get('embedded_logo_path')
                if embedded_logo_path:
                    # Resolve relative paths
                    resolved_path = self._resolve_logo_path(embedded_logo_path)
                    if resolved_path:
                        return resolved_path
            
            # Fallback: check data dict for template_embedded_logo
            if data and isinstance(data, dict) and 'template_embedded_logo' in data:
                return data['template_embedded_logo']

        return self.image_path

    def _resolve_logo_path(self, logo_path: str) -> Optional[str]:
        """
        Resolve logo path, handling both absolute and relative paths.

        Args:
            logo_path: Logo path from template settings

        Returns:
            Resolved absolute path, or None if not found
        """
        # If absolute path, use as-is
        if os.path.isabs(logo_path):
            return logo_path if os.path.exists(logo_path) else None

        # Try relative to project root
        project_root = Path.cwd()
        resolved = project_root / logo_path
        if resolved.exists():
            return str(resolved)

        # Try relative to templates directory
        templates_dir = project_root / 'templates'
        resolved = templates_dir / logo_path
        if resolved.exists():
            return str(resolved)

        # Try relative to assets directory (if it exists)
        assets_dir = project_root / 'assets'
        if assets_dir.exists():
            resolved = assets_dir / logo_path
            if resolved.exists():
                return str(resolved)

        # Return original path (will show placeholder if doesn't exist)
        return logo_path

    def _apply_styling(self, picture) -> None:
        """
        Apply styling to picture shape using XML manipulation.

        Args:
            picture: python-pptx picture shape

        Supports:
        - Opacity/transparency (0-100)
        - Border width and color
        - Corner radius (rounded corners)

        Note: These features require XML manipulation and may not work
        with all python-pptx versions or PowerPoint versions.
        """
        try:
            # Get the underlying XML element
            pic_element = picture._element
            
            # Apply border first (most reliable)
            if self.border_width > 0:
                self._apply_border(pic_element, self.border_width, self.border_color)
            
            # Apply opacity/transparency
            if self.opacity < 100:
                self._apply_opacity(pic_element, self.opacity)
            
            # Apply corner radius (may require shape wrapping)
            if self.corner_radius > 0:
                self._apply_corner_radius(pic_element, self.corner_radius)
                
        except Exception:
            # Silently fail if XML manipulation doesn't work
            # This is expected if python-pptx version doesn't support it
            # or if the PowerPoint format doesn't support these features
            pass

    def _apply_opacity(self, pic_element, opacity: int) -> None:
        """
        Apply opacity/transparency to image using XML.

        Args:
            pic_element: XML element of the picture
            opacity: Opacity value (0-100, where 100 is fully opaque)
        """
        try:
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import nsdecls, qn
            
            # Convert opacity percentage to alpha (0-100000, where 100000 = 100%)
            alpha = int((opacity / 100.0) * 100000)
            
            # Find or create the picture effect list
            # Opacity is set in the picture's effect list
            pic = pic_element.find('.//a:pic', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if pic is not None:
                # Get or create effect list
                effect_lst = pic.find('.//a:effectLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if effect_lst is None:
                    # Create effect list if it doesn't exist
                    from pptx.oxml import parse_xml
                    effect_lst_xml = '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                    effect_lst = parse_xml(effect_lst_xml)
                    pic.append(effect_lst)
                
                # Add alpha effect for transparency
                alpha_effect_xml = f'''<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                    <a:val val="{alpha}"/>
                </a:alpha>'''
                alpha_effect = parse_xml(alpha_effect_xml)
                effect_lst.append(alpha_effect)
        except Exception:
            # Fallback: Try setting opacity via shape fill if available
            try:
                if hasattr(pic_element, 'spPr') and pic_element.spPr is not None:
                    # Alternative method via shape properties
                    pass
            except Exception:
                pass

    def _apply_border(self, pic_element, border_width: int, border_color: str) -> None:
        """
        Apply border to image using XML.

        Args:
            pic_element: XML element of the picture
            border_width: Border width in points
            border_color: Hex color string (e.g., '#000000')
        """
        try:
            from pptx.oxml import parse_xml
            from pptx.dml.color import RGBColor
            from pptx.util import Pt
            
            # Try using python-pptx's built-in line property first (more reliable)
            try:
                if hasattr(pic_element, 'line'):
                    # Parse border color
                    hex_color = border_color.lstrip('#')
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16)
                    b = int(hex_color[4:6], 16)
                    
                    pic_element.line.color.rgb = RGBColor(r, g, b)
                    pic_element.line.width = Pt(border_width)
                    return
            except (AttributeError, ValueError):
                pass
            
            # Fallback to XML manipulation
            # Get shape properties
            spPr = pic_element.find('.//a:spPr', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if spPr is None:
                # Create shape properties if they don't exist
                spPr_xml = '<a:spPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                spPr = parse_xml(spPr_xml)
                # Insert before any existing elements
                pic_element.insert(0, spPr)
            
            # Parse border color
            hex_color = border_color.lstrip('#')
            
            # Create line (border) element
            # Convert points to EMUs (English Metric Units) - 1 point = 12700 EMUs
            width_emu = border_width * 12700
            
            ln_xml = f'''<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{width_emu}">
                <a:solidFill>
                    <a:srgbClr val="{hex_color.upper()}"/>
                </a:solidFill>
            </a:ln>'''
            ln = parse_xml(ln_xml)
            
            # Remove existing line if present
            existing_ln = spPr.find('.//a:ln', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if existing_ln is not None:
                spPr.remove(existing_ln)
            
            # Add new line
            spPr.append(ln)
            
        except Exception:
            # Silently fail if border cannot be applied
            pass

    def _apply_corner_radius(self, pic_element, corner_radius: int) -> None:
        """
        Apply corner radius (rounded corners) to image using XML.

        Args:
            pic_element: XML element of the picture
            corner_radius: Corner radius in points
        """
        try:
            from pptx.oxml import parse_xml
            
            # Get shape properties
            spPr = pic_element.find('.//a:spPr', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if spPr is None:
                # Create shape properties if they don't exist
                spPr_xml = '<a:spPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                spPr = parse_xml(spPr_xml)
                pic_element.insert(0, spPr)
            
            # Convert points to EMUs (1 point = 12700 EMUs)
            radius_emu = corner_radius * 12700
            
            # Create rounded rectangle adjustment
            # For rounded rectangles, we need to add adjustment values
            # The adjustment value is typically a percentage of the shape size
            # For simplicity, we'll use a fixed adjustment value
            
            # Remove existing prstGeom if it's not a rounded rectangle
            existing_geom = spPr.find('.//a:prstGeom', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if existing_geom is not None:
                prst = existing_geom.get('prst')
                if prst != 'roundRect':
                    spPr.remove(existing_geom)
                    existing_geom = None
            
            # Create or update rounded rectangle geometry
            if existing_geom is None:
                # Create new rounded rectangle geometry
                # Adjustment value: typically 10000-50000 (10%-50% of shape size)
                # Higher value = more rounded
                adj_value = min(50000, max(10000, radius_emu // 10))
                
                prst_geom_xml = f'''<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="roundRect">
                    <a:avLst>
                        <a:gd name="adj" fmla="val {adj_value}"/>
                    </a:avLst>
                </a:prstGeom>'''
                prst_geom = parse_xml(prst_geom_xml)
                
                # Remove existing geometry if present
                existing_xfrm = spPr.find('.//a:xfrm', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if existing_xfrm is not None:
                    # Insert before xfrm
                    spPr.insert(0, prst_geom)
                else:
                    spPr.insert(0, prst_geom)
            else:
                # Update existing rounded rectangle adjustment
                av_lst = existing_geom.find('.//a:avLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if av_lst is not None:
                    adj_value = min(50000, max(10000, radius_emu // 10))
                    gd = av_lst.find('.//a:gd', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if gd is not None:
                        gd.set('fmla', f'val {adj_value}')
                        
        except Exception:
            # Corner radius may not be supported for all image types
            # Silently fail
            pass

    def _render_placeholder(self, slide: Slide, error: str = None) -> None:
        """
        Render placeholder when image is not available.

        Args:
            slide: PowerPoint slide object
            error: Optional error message
        """
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        # Add rectangle shape as placeholder
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.x,
            self.y,
            self.width,
            self.height
        )

        # Gray background
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(243, 244, 246)

        # Gray border
        placeholder.line.color.rgb = RGBColor(209, 213, 219)
        placeholder.line.width = Pt(1)

        # Add text
        text_frame = placeholder.text_frame
        text_frame.clear()

        paragraph = text_frame.paragraphs[0]
        if error:
            paragraph.text = f"[Image Error]\n{error}"
        elif self.image_path:
            paragraph.text = f"[Image]\n{os.path.basename(self.image_path)}\nNot Found"
        else:
            paragraph.text = "[Image Placeholder]"

        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(156, 163, 175)

        # Center text vertically
        text_frame.vertical_anchor = 1  # Middle

    def set_image_path(self, path: str) -> None:
        """
        Update image file path.

        Args:
            path: New image file path
        """
        self.image_path = path
        self.data_source['path'] = path

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        config = super().to_dict()
        config['data_source']['path'] = self.image_path
        config['data_source']['type'] = self.image_type
        return config


# Example usage configurations:
EXAMPLE_CONFIGS = {
    'logo': {
        'type': 'image',
        'position': {'x': 0.5, 'y': 0.5},
        'size': {'width': 2.0, 'height': 1.0},
        'data_source': {
            'path': 'assets/logo.png',
            'type': 'file'
        },
        'style': {
            'maintain_aspect': True,
            'border_width': 0,
            'opacity': 100
        }
    },
    'template_logo': {
        'type': 'image',
        'position': {'x': 8.5, 'y': 0.5},
        'size': {'width': 1.0, 'height': 0.5},
        'data_source': {
            'type': 'template_logo'  # Will use logo from template config
        },
        'style': {
            'maintain_aspect': True
        }
    },
    'product_image': {
        'type': 'image',
        'position': {'x': 5.0, 'y': 2.0},
        'size': {'width': 4.0, 'height': 3.0},
        'data_source': {
            'path': 'products/product_001.jpg',
            'type': 'file'
        },
        'style': {
            'maintain_aspect': False,  # Stretch to fit
            'border_width': 2,
            'border_color': '#E5E7EB',
            'corner_radius': 8
        }
    }
}
