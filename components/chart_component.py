"""
ChartComponent - Renders data visualizations using matplotlib

Used for: Bar charts, column charts, pie charts, line charts, trend analysis
Supports: Multiple chart types, custom colors, legends, axes labels
"""

from typing import Any, Dict, List, Optional
from pptx.slide import Slide
from components.base_component import BaseComponent
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib
import os
import tempfile
import warnings

# Use non-interactive backend for server environments
matplotlib.use('Agg')

# Suppress matplotlib tight_layout warnings
warnings.filterwarnings('ignore', message='.*Tight layout.*')


class ChartComponent(BaseComponent):
    """
    Component for rendering charts and data visualizations.

    Supports: Column, Bar, Pie, Line, Stacked Column, Stacked Bar charts
    Uses matplotlib for chart generation, then embeds as image in PowerPoint.
    """

    # Chart type constants
    CHART_COLUMN = 'column'
    CHART_BAR = 'bar'
    CHART_PIE = 'pie'
    CHART_LINE = 'line'
    CHART_STACKED_COLUMN = 'stacked_column'
    CHART_STACKED_BAR = 'stacked_bar'

    def __init__(self, config: Dict[str, Any]):
        """
        Initialize ChartComponent.

        Config keys:
            - chart_type: 'column' | 'bar' | 'pie' | 'line' | 'stacked_column' | 'stacked_bar'
            - data_source:
                - x_column: Column for X-axis (categories)
                - y_column: Column for Y-axis (values)
                - series_column: Column for series (multi-series charts)
                - sort_by: Column to sort by
                - top_n: Show only top N items
            - style:
                - colors: List of hex colors or 'brand' for template colors
                - legend_position: 'top' | 'bottom' | 'left' | 'right' | 'none'
                - show_values: Show data labels on chart
                - title: Chart title (optional)
                - x_label: X-axis label
                - y_label: Y-axis label
                - grid: Show grid lines
                - font_size: Chart font size
        """
        super().__init__(config)

        # Chart type
        self.chart_type = config.get('chart_type', self.CHART_COLUMN)

        # Data configuration
        self.x_column = self.data_source.get('x_column')
        self.y_column = self.data_source.get('y_column')
        self.series_column = self.data_source.get('series_column')
        self.sort_by = self.data_source.get('sort_by')
        self.top_n = self.data_source.get('top_n')

        # Chart styling
        self.colors = self.style.get('colors', 'default')
        self.legend_position = self.style.get('legend_position', 'bottom')
        self.show_values = self.style.get('show_values', False)
        self.chart_title = self.style.get('title', '')
        self.x_label = self.style.get('x_label', '')
        self.y_label = self.style.get('y_label', '')
        self.show_grid = self.style.get('grid', True)
        self.chart_font_size = self.style.get('font_size', 10)

        # Validate after all attributes are set
        self.validate()

    def validate(self) -> bool:
        """Validate chart component configuration."""
        super().validate()

        valid_types = [
            self.CHART_COLUMN, self.CHART_BAR, self.CHART_PIE,
            self.CHART_LINE, self.CHART_STACKED_COLUMN, self.CHART_STACKED_BAR
        ]
        if self.chart_type not in valid_types:
            raise ValueError(f"Invalid chart_type: {self.chart_type}. Must be one of {valid_types}")

        if not self.x_column and self.chart_type != self.CHART_PIE:
            raise ValueError("x_column is required for non-pie charts")

        if not self.y_column:
            raise ValueError("y_column is required")

        return True

    def render(self, slide: Slide, data: Any) -> None:
        """
        Render chart on the slide.

        Args:
            slide: PowerPoint slide object
            data: pandas DataFrame or similar data structure
        """
        # Prepare data
        df = self._prepare_data(data)

        if df is None or df.empty:
            self._render_placeholder(slide)
            return

        # Generate chart image
        chart_image_path = self._generate_chart(df)

        if not chart_image_path:
            self._render_placeholder(slide)
            return

        try:
            # Add chart image to slide
            slide.shapes.add_picture(
                chart_image_path,
                self.x,
                self.y,
                width=self.width,
                height=self.height
            )
        finally:
            # Clean up temporary file
            if os.path.exists(chart_image_path):
                try:
                    os.unlink(chart_image_path)
                except OSError:
                    pass  # File already deleted or doesn't exist

    def _prepare_data(self, data: Any) -> Optional[pd.DataFrame]:
        """
        Convert and prepare data for charting.

        Args:
            data: Input data (DataFrame, list, dict)

        Returns:
            pd.DataFrame or None
        """
        if data is None:
            return None

        # Convert to DataFrame
        if isinstance(data, pd.DataFrame):
            df = data.copy()
        elif isinstance(data, list):
            if not data:
                return None
            df = pd.DataFrame(data)
        elif isinstance(data, dict):
            df = pd.DataFrame(data)
        else:
            return None

        # Filter columns if specified
        required_cols = [self.x_column, self.y_column]
        if self.series_column:
            required_cols.append(self.series_column)

        available_cols = [col for col in required_cols if col and col in df.columns]
        if not available_cols:
            return None

        # Drop rows with NaN in critical columns
        if self.y_column and self.y_column in df.columns:
            df = df.dropna(subset=[self.y_column])

        if self.x_column and self.x_column in df.columns:
            df = df.dropna(subset=[self.x_column])

        # Check if we have any data left
        if df.empty:
            return None

        # Apply sorting
        if self.sort_by and self.sort_by in df.columns:
            df = df.sort_values(by=self.sort_by, ascending=False)

        # Apply top N filter
        if self.top_n and self.top_n > 0:
            df = df.head(self.top_n)

        return df

    def _generate_chart(self, df: pd.DataFrame) -> Optional[str]:
        """
        Generate chart using matplotlib and save to temp file.

        Args:
            df: DataFrame with chart data

        Returns:
            str: Path to temporary image file
        """
        try:
            # Set figure size based on component dimensions (limit to reasonable values)
            # Ensure we have valid numeric values
            try:
                raw_width = float(self.width.inches)
                raw_height = float(self.height.inches)
            except (ValueError, AttributeError, TypeError):
                raw_width = 9.0
                raw_height = 5.0

            # DEBUG: Print values to help diagnose
            print(f"[DEBUG] Raw dimensions: {raw_width}x{raw_height} inches")

            # Clamp to reasonable bounds - FORCE to safe values
            if raw_width > 20.0 or raw_width < 1.0:
                print(f"[DEBUG] Clamping width from {raw_width} to 20.0")
                fig_width = 20.0 if raw_width > 20.0 else 9.0
            else:
                fig_width = raw_width

            if raw_height > 15.0 or raw_height < 1.0:
                print(f"[DEBUG] Clamping height from {raw_height} to 15.0")
                fig_height = 15.0 if raw_height > 15.0 else 5.0
            else:
                fig_height = raw_height

            # Use lower DPI to prevent huge images
            dpi = 100  # Reduced from 150

            print(f"[DEBUG] Final dimensions: {fig_width}x{fig_height} inches at {dpi} DPI = {fig_width*dpi}x{fig_height*dpi} pixels")

            # Check if resulting image would be too large (max 2^16 = 65536)
            max_pixels = 10000  # Conservative limit
            if fig_width * dpi > max_pixels or fig_height * dpi > max_pixels:
                # Scale down DPI to fit within limits
                old_dpi = dpi
                dpi = int(min(max_pixels / fig_width, max_pixels / fig_height, dpi))
                dpi = max(50, dpi)  # Minimum 50 DPI for readability
                print(f"[DEBUG] Scaled DPI from {old_dpi} to {dpi}")

            # Set matplotlib's default figsize and dpi to prevent pandas .plot() from creating huge figures
            # This is critical because pandas DataFrame.plot() creates figures internally using rcParams
            plt.rcParams['figure.figsize'] = [fig_width, fig_height]
            plt.rcParams['figure.dpi'] = dpi
            plt.rcParams['figure.max_open_warning'] = 0  # Suppress warnings about too many figures

            fig, ax = plt.subplots(figsize=(fig_width, fig_height), dpi=dpi)

            # Get colors
            colors = self._get_colors(df)

            # Generate appropriate chart type
            if self.chart_type == self.CHART_COLUMN:
                self._create_column_chart(ax, df, colors)
            elif self.chart_type == self.CHART_BAR:
                self._create_bar_chart(ax, df, colors)
            elif self.chart_type == self.CHART_PIE:
                self._create_pie_chart(ax, df, colors)
            elif self.chart_type == self.CHART_LINE:
                self._create_line_chart(ax, df, colors)
            elif self.chart_type == self.CHART_STACKED_COLUMN:
                self._create_stacked_column_chart(ax, df, colors)
            elif self.chart_type == self.CHART_STACKED_BAR:
                self._create_stacked_bar_chart(ax, df, colors)

            # Apply common styling
            self._apply_chart_styling(ax, fig)

            # Save to temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            temp_path = temp_file.name
            temp_file.close()

            # Apply tight layout (suppress warnings if it fails)
            try:
                plt.tight_layout(pad=1.0)  # Use padding instead of tight bbox
            except Exception:
                pass  # Ignore layout warnings

            # Calculate expected pixel dimensions
            expected_width_px = int(fig_width * dpi)
            expected_height_px = int(fig_height * dpi)
            
            # Maximum allowed dimensions (2^16 = 65536, use 10000 for safety)
            MAX_DIMENSION_PX = 10000
            
            # Validate dimensions before saving
            if expected_width_px > MAX_DIMENSION_PX or expected_height_px > MAX_DIMENSION_PX:
                # Scale down if still too large
                scale_factor = min(MAX_DIMENSION_PX / expected_width_px, MAX_DIMENSION_PX / expected_height_px, 1.0)
                new_dpi = max(50, int(dpi * scale_factor))
                print(f"[DEBUG] Scaling DPI from {dpi} to {new_dpi} to fit within {MAX_DIMENSION_PX}px limit")
                dpi = new_dpi
                expected_width_px = int(fig_width * dpi)
                expected_height_px = int(fig_height * dpi)

            # Save with validated DPI, use pad_inches instead of bbox_inches='tight'
            # bbox_inches='tight' can cause images to exceed expected dimensions
            plt.savefig(
                temp_path, 
                dpi=dpi, 
                bbox_inches=None,  # Don't use 'tight' - it causes dimension issues
                pad_inches=0.1,  # Small padding instead
                facecolor='white',
                edgecolor='none'
            )
            
            # Verify saved image dimensions (check actual file size if possible)
            try:
                from PIL import Image
                with Image.open(temp_path) as img:
                    actual_width, actual_height = img.size
                    if actual_width > MAX_DIMENSION_PX or actual_height > MAX_DIMENSION_PX:
                        print(f"[WARNING] Saved image is {actual_width}x{actual_height}px, exceeds {MAX_DIMENSION_PX}px limit")
                        # Resize if needed
                        if actual_width > MAX_DIMENSION_PX:
                            scale = MAX_DIMENSION_PX / actual_width
                            new_width = MAX_DIMENSION_PX
                            new_height = int(actual_height * scale)
                            img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                            img_resized.save(temp_path)
                            print(f"[DEBUG] Resized image to {new_width}x{new_height}px")
            except ImportError:
                pass  # PIL not available, skip verification
            except Exception as e:
                print(f"[DEBUG] Could not verify image size: {e}")
            
            plt.close(fig)

            return temp_path

        except Exception as e:
            print(f"Error generating chart: {e}")
            import traceback
            traceback.print_exc()  # Print full traceback for debugging
            if 'fig' in locals():
                try:
                    plt.close(fig)
                except Exception:
                    pass
            return None

    def _create_column_chart(self, ax, df: pd.DataFrame, colors: List[str]) -> None:
        """Create vertical column chart."""
        if self.series_column and self.series_column in df.columns:
            # Multi-series column chart
            df_pivot = df.pivot_table(
                values=self.y_column,
                index=self.x_column,
                columns=self.series_column,
                aggfunc='sum'
            )
            df_pivot.plot(kind='bar', ax=ax, color=colors, width=0.7)
        else:
            # Single series
            x_data = df[self.x_column].astype(str)
            y_data = df[self.y_column]
            ax.bar(x_data, y_data, color=colors[0] if colors else None, width=0.6)

            if self.show_values:
                for i, v in enumerate(y_data):
                    ax.text(i, v, f'{v:,.0f}', ha='center', va='bottom', fontsize=8)

    def _create_bar_chart(self, ax, df: pd.DataFrame, colors: List[str]) -> None:
        """Create horizontal bar chart."""
        if self.series_column and self.series_column in df.columns:
            # Multi-series bar chart
            df_pivot = df.pivot_table(
                values=self.y_column,
                index=self.x_column,
                columns=self.series_column,
                aggfunc='sum'
            )
            df_pivot.plot(kind='barh', ax=ax, color=colors, height=0.7)
        else:
            # Single series
            x_data = df[self.x_column].astype(str)
            y_data = df[self.y_column]
            ax.barh(x_data, y_data, color=colors[0] if colors else None, height=0.6)

            if self.show_values:
                for i, v in enumerate(y_data):
                    ax.text(v, i, f'{v:,.0f}', ha='left', va='center', fontsize=8)

    def _create_pie_chart(self, ax, df: pd.DataFrame, colors: List[str]) -> None:
        """Create pie chart."""
        labels = df[self.x_column].astype(str) if self.x_column else df.index
        values = df[self.y_column]

        wedges, texts, autotexts = ax.pie(
            values,
            labels=labels,
            colors=colors,
            autopct='%1.1f%%' if self.show_values else None,
            startangle=90
        )

        # Style percentages
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontsize(9)
            autotext.set_weight('bold')

        ax.axis('equal')

    def _create_line_chart(self, ax, df: pd.DataFrame, colors: List[str]) -> None:
        """Create line chart."""
        if self.series_column and self.series_column in df.columns:
            # Multi-series line chart
            for idx, series_name in enumerate(df[self.series_column].unique()):
                series_data = df[df[self.series_column] == series_name]
                color = colors[idx % len(colors)] if colors else None
                ax.plot(
                    series_data[self.x_column],
                    series_data[self.y_column],
                    marker='o',
                    label=series_name,
                    color=color,
                    linewidth=2
                )
        else:
            # Single series
            ax.plot(
                df[self.x_column],
                df[self.y_column],
                marker='o',
                color=colors[0] if colors else None,
                linewidth=2
            )

            if self.show_values:
                for x, y in zip(df[self.x_column], df[self.y_column]):
                    ax.text(x, y, f'{y:,.0f}', ha='center', va='bottom', fontsize=8)

    def _create_stacked_column_chart(self, ax, df: pd.DataFrame, colors: List[str]) -> None:
        """Create stacked vertical column chart."""
        if not self.series_column or self.series_column not in df.columns:
            # Fallback to regular column
            self._create_column_chart(ax, df, colors)
            return

        df_pivot = df.pivot_table(
            values=self.y_column,
            index=self.x_column,
            columns=self.series_column,
            aggfunc='sum'
        )
        df_pivot.plot(kind='bar', stacked=True, ax=ax, color=colors, width=0.7)

    def _create_stacked_bar_chart(self, ax, df: pd.DataFrame, colors: List[str]) -> None:
        """Create stacked horizontal bar chart."""
        if not self.series_column or self.series_column not in df.columns:
            # Fallback to regular bar
            self._create_bar_chart(ax, df, colors)
            return

        df_pivot = df.pivot_table(
            values=self.y_column,
            index=self.x_column,
            columns=self.series_column,
            aggfunc='sum'
        )
        df_pivot.plot(kind='barh', stacked=True, ax=ax, color=colors, height=0.7)

    def _get_colors(self, df: pd.DataFrame) -> List[str]:
        """
        Get color list for chart.

        Args:
            df: DataFrame with chart data

        Returns:
            List of hex color codes (always returns at least one color)
        """
        # Try user-defined colors first
        if isinstance(self.colors, list) and len(self.colors) > 0:
            return self.colors

        if self.colors == 'brand':
            # Use template brand colors (TODO: get from template)
            return ['#2563EB', '#10B981', '#F59E0B', '#EF4444', '#8B5CF6', '#EC4899']

        # Default matplotlib colors - with robust error handling
        try:
            prop_cycle = plt.rcParams['axes.prop_cycle']
            # Handle different prop_cycle formats safely
            colors = []
            for c in prop_cycle:
                try:
                    if isinstance(c, dict):
                        color = c.get('color') or c.get('c')
                    elif hasattr(c, 'color'):
                        color = c.color
                    elif hasattr(c, 'c'):
                        color = c.c
                    else:
                        continue
                    if color:
                        colors.append(color)
                except (KeyError, AttributeError, IndexError, TypeError):
                    continue
                if len(colors) >= 10:
                    break
            
            if colors:
                return colors
        except (KeyError, AttributeError, IndexError, TypeError) as e:
            # Fall through to default colors
            pass

        # Fallback colors - always return at least one color
        return ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd', '#8c564b']

    def _apply_chart_styling(self, ax, fig) -> None:
        """
        Apply styling to chart (title, labels, grid, legend).

        Args:
            ax: Matplotlib axes
            fig: Matplotlib figure
        """
        # Title
        if self.chart_title:
            ax.set_title(self.chart_title, fontsize=self.chart_font_size + 2, fontweight='bold')

        # Axis labels
        if self.x_label and self.chart_type != self.CHART_PIE:
            ax.set_xlabel(self.x_label, fontsize=self.chart_font_size)
        if self.y_label and self.chart_type != self.CHART_PIE:
            ax.set_ylabel(self.y_label, fontsize=self.chart_font_size)

        # Grid
        if self.show_grid and self.chart_type != self.CHART_PIE:
            ax.grid(True, alpha=0.3, linestyle='--', linewidth=0.5)
            ax.set_axisbelow(True)

        # Legend
        if self.legend_position != 'none' and self.series_column:
            loc_map = {
                'top': 'upper center',
                'bottom': 'lower center',
                'left': 'center left',
                'right': 'center right'
            }
            ax.legend(
                loc=loc_map.get(self.legend_position, 'best'),
                fontsize=self.chart_font_size - 1,
                frameon=False
            )

        # Font sizes
        ax.tick_params(axis='both', labelsize=self.chart_font_size - 1)

        # Remove top and right spines for cleaner look
        if self.chart_type not in [self.CHART_PIE]:
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)

    def _render_placeholder(self, slide: Slide) -> None:
        """Render placeholder when chart cannot be generated."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.x, self.y, self.width, self.height
        )

        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(243, 244, 246)
        placeholder.line.color.rgb = RGBColor(209, 213, 219)

        text_frame = placeholder.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.text = f"[{self.chart_type.upper()} Chart]\nNo data available"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(156, 163, 175)
        text_frame.vertical_anchor = 1

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        config = super().to_dict()
        config['chart_type'] = self.chart_type
        config['data_source'].update({
            'x_column': self.x_column,
            'y_column': self.y_column,
            'series_column': self.series_column,
            'sort_by': self.sort_by,
            'top_n': self.top_n
        })
        return config


# Example configurations:
EXAMPLE_CONFIGS = {
    'column_chart': {
        'type': 'chart',
        'chart_type': 'column',
        'position': {'x': 0.5, 'y': 2.0},
        'size': {'width': 9.0, 'height': 4.0},
        'data_source': {
            'x_column': 'Company',
            'y_column': 'Total',
            'sort_by': 'Total',
            'top_n': 10
        },
        'style': {
            'colors': ['#2563EB', '#10B981', '#F59E0B'],
            'legend_position': 'none',
            'show_values': True,
            'title': 'Media Mentions by Company',
            'y_label': 'Number of Mentions',
            'grid': True
        }
    },
    'pie_chart': {
        'type': 'chart',
        'chart_type': 'pie',
        'position': {'x': 5.0, 'y': 2.0},
        'size': {'width': 4.5, 'height': 4.0},
        'data_source': {
            'x_column': 'Category',
            'y_column': 'Percentage'
        },
        'style': {
            'colors': ['#2563EB', '#10B981', '#F59E0B', '#EF4444'],
            'show_values': True,
            'title': 'Sentiment Distribution'
        }
    },
    'line_chart': {
        'type': 'chart',
        'chart_type': 'line',
        'position': {'x': 0.5, 'y': 2.0},
        'size': {'width': 9.0, 'height': 4.0},
        'data_source': {
            'x_column': 'Month',
            'y_column': 'Mentions',
            'series_column': 'Company'
        },
        'style': {
            'colors': ['#2563EB', '#10B981', '#F59E0B'],
            'legend_position': 'bottom',
            'show_values': False,
            'title': '12-Month Trend Analysis',
            'x_label': 'Month',
            'y_label': 'Media Mentions',
            'grid': True
        }
    }
}
