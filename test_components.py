"""
Test script for all ReportForge components

Run this to verify all components are working correctly.
Generates a test PowerPoint file with all component types.
"""

from pptx import Presentation
from pptx.util import Inches
from components import (
    TextComponent,
    TableComponent,
    ImageComponent,
    ChartComponent,
    SummaryComponent
)
import pandas as pd
from datetime import datetime


def create_sample_data():
    """Create sample data for testing."""
    return pd.DataFrame({
        'Company': ['BSH', 'Arçelik', 'Vestel', 'Beko', 'Profilo'],
        'Total': [1234, 987, 654, 543, 321],
        'Positive': [856, 654, 432, 378, 210],
        'Negative': [234, 187, 123, 98, 67],
        'Neutral': [144, 146, 99, 67, 44]
    })


def create_trend_data():
    """Create trend data for time-series testing."""
    return pd.DataFrame({
        'Month': ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'],
        'BSH': [100, 120, 115, 140, 160, 175],
        'Arçelik': [90, 95, 100, 105, 110, 115],
        'Vestel': [70, 75, 72, 80, 85, 90]
    })


def test_text_component(slide):
    """Test TextComponent."""
    print("Testing TextComponent...")

    # Title
    title = TextComponent({
        'type': 'text',
        'content': 'ReportForge Component Test - {date}',
        'position': {'x': 0.5, 'y': 0.5},
        'size': {'width': 9.0, 'height': 1.0},
        'variables': {'date': datetime.now().strftime('%Y-%m-%d')},
        'style': {
            'font_name': 'Calibri',
            'font_size': 32,
            'bold': True,
            'color': '#1F2937',
            'alignment': 'center'
        }
    })
    title.render(slide)

    # Subtitle
    subtitle = TextComponent({
        'type': 'text',
        'content': 'Testing all component types: Text, Table, Image, Chart, Summary',
        'position': {'x': 0.5, 'y': 1.5},
        'size': {'width': 9.0, 'height': 0.5},
        'style': {
            'font_size': 14,
            'color': '#6B7280',
            'alignment': 'center'
        }
    })
    subtitle.render(slide)

    print("✅ TextComponent test complete")


def test_table_component(slide, data):
    """Test TableComponent."""
    print("Testing TableComponent...")

    table = TableComponent({
        'type': 'table',
        'position': {'x': 0.5, 'y': 2.5},
        'size': {'width': 9.0, 'height': 3.0},
        'data_source': {
            'columns': ['Company', 'Total', 'Positive', 'Negative', 'Neutral'],
            'sort_by': 'Total',
            'ascending': False
        },
        'style': {
            'header_row': True,
            'zebra_striping': True,
            'border': True,
            'header_color': '#2563EB',
            'header_text_color': '#FFFFFF',
            'row_color_1': '#FFFFFF',
            'row_color_2': '#F9FAFB',
            'text_color': '#1F2937',
            'font_name': 'Calibri',
            'font_size': 10
        }
    })
    table.render(slide, data)

    print("✅ TableComponent test complete")


def test_image_component(slide):
    """Test ImageComponent."""
    print("Testing ImageComponent...")

    # Note: This will show a placeholder since we don't have an actual image
    image = ImageComponent({
        'type': 'image',
        'position': {'x': 8.5, 'y': 0.5},
        'size': {'width': 1.0, 'height': 0.8},
        'data_source': {
            'path': 'assets/test_logo.png',  # Will show placeholder if missing
            'type': 'file'
        },
        'style': {
            'maintain_aspect': True
        }
    })
    image.render(slide)

    print("✅ ImageComponent test complete (placeholder shown)")


def test_chart_component_slide(prs, data):
    """Test ChartComponent with multiple chart types."""
    print("Testing ChartComponent...")

    # Slide 2: Column Chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = TextComponent({
        'type': 'text',
        'content': 'Chart Component Tests',
        'position': {'x': 0.5, 'y': 0.3},
        'size': {'width': 9.0, 'height': 0.7},
        'style': {'font_size': 28, 'bold': True, 'alignment': 'center'}
    })
    title.render(slide)

    # Column chart
    column_chart = ChartComponent({
        'type': 'chart',
        'chart_type': 'column',
        'position': {'x': 0.5, 'y': 1.2},
        'size': {'width': 4.5, 'height': 3.0},
        'data_source': {
            'x_column': 'Company',
            'y_column': 'Total',
            'sort_by': 'Total',
            'top_n': 5
        },
        'style': {
            'colors': ['#2563EB'],
            'show_values': True,
            'title': 'Total Mentions by Company',
            'y_label': 'Mentions',
            'grid': True,
            'legend_position': 'none'
        }
    })
    column_chart.render(slide, data)

    # Pie chart
    pie_data = pd.DataFrame({
        'Category': ['Positive', 'Negative', 'Neutral'],
        'Count': [data['Positive'].sum(), data['Negative'].sum(), data['Neutral'].sum()]
    })

    pie_chart = ChartComponent({
        'type': 'chart',
        'chart_type': 'pie',
        'position': {'x': 5.2, 'y': 1.2},
        'size': {'width': 4.3, 'height': 3.0},
        'data_source': {
            'x_column': 'Category',
            'y_column': 'Count'
        },
        'style': {
            'colors': ['#10B981', '#EF4444', '#F59E0B'],
            'show_values': True,
            'title': 'Sentiment Distribution'
        }
    })
    pie_chart.render(slide, pie_data)

    # Bar chart (horizontal)
    bar_chart = ChartComponent({
        'type': 'chart',
        'chart_type': 'bar',
        'position': {'x': 0.5, 'y': 4.5},
        'size': {'width': 9.0, 'height': 2.5},
        'data_source': {
            'x_column': 'Company',
            'y_column': 'Positive',
            'sort_by': 'Positive'
        },
        'style': {
            'colors': ['#10B981'],
            'show_values': True,
            'title': 'Positive Mentions (Bar Chart)',
            'x_label': 'Number of Mentions',
            'grid': True
        }
    })
    bar_chart.render(slide, data)

    print("✅ ChartComponent test complete (Column, Pie, Bar)")


def test_line_chart_slide(prs, trend_data):
    """Test line chart with trend data."""
    print("Testing LineChart...")

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = TextComponent({
        'type': 'text',
        'content': 'Line Chart Test - Trend Analysis',
        'position': {'x': 0.5, 'y': 0.3},
        'size': {'width': 9.0, 'height': 0.7},
        'style': {'font_size': 28, 'bold': True, 'alignment': 'center'}
    })
    title.render(slide)

    # Prepare data for multi-series line chart
    line_data = pd.melt(
        trend_data,
        id_vars=['Month'],
        value_vars=['BSH', 'Arçelik', 'Vestel'],
        var_name='Company',
        value_name='Mentions'
    )

    line_chart = ChartComponent({
        'type': 'chart',
        'chart_type': 'line',
        'position': {'x': 0.5, 'y': 1.2},
        'size': {'width': 9.0, 'height': 5.0},
        'data_source': {
            'x_column': 'Month',
            'y_column': 'Mentions',
            'series_column': 'Company'
        },
        'style': {
            'colors': ['#2563EB', '#10B981', '#F59E0B'],
            'legend_position': 'bottom',
            'title': '6-Month Trend Analysis',
            'x_label': 'Month',
            'y_label': 'Media Mentions',
            'grid': True
        }
    })
    line_chart.render(slide, line_data)

    print("✅ LineChart test complete")


def test_summary_component_slide(prs, data, trend_data):
    """Test SummaryComponent."""
    print("Testing SummaryComponent...")

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = TextComponent({
        'type': 'text',
        'content': 'Summary Component Test - Auto-Generated Insights',
        'position': {'x': 0.5, 'y': 0.3},
        'size': {'width': 9.0, 'height': 0.7},
        'style': {'font_size': 28, 'bold': True, 'alignment': 'center'}
    })
    title.render(slide)

    # Key metrics summary
    summary1 = SummaryComponent({
        'type': 'summary',
        'position': {'x': 0.5, 'y': 1.2},
        'size': {'width': 9.0, 'height': 1.5},
        'data_source': {
            'insight_types': ['key_metrics'],
            'metric_columns': ['Total', 'Positive', 'Negative'],
            'max_items': 3
        },
        'style': {
            'layout': 'bullets',
            'show_icons': True,
            'font_size': 14
        }
    })
    summary1.render(slide, data)

    # Highlights and comparisons
    summary2 = SummaryComponent({
        'type': 'summary',
        'position': {'x': 0.5, 'y': 3.0},
        'size': {'width': 9.0, 'height': 2.0},
        'data_source': {
            'insight_types': ['highlights', 'comparisons'],
            'metric_columns': ['Total', 'Positive'],
            'compare_column': 'Company',
            'max_items': 4
        },
        'style': {
            'layout': 'callout_boxes',
            'show_icons': True,
            'highlight_color': '#EFF6FF'
        }
    })
    summary2.render(slide, data)

    # Top performers
    summary3 = SummaryComponent({
        'type': 'summary',
        'position': {'x': 0.5, 'y': 5.3},
        'size': {'width': 9.0, 'height': 1.5},
        'data_source': {
            'insight_types': ['top_performers'],
            'metric_columns': ['Total', 'Positive'],
            'compare_column': 'Company',
            'max_items': 2
        },
        'style': {
            'layout': 'numbered',
            'show_icons': True
        }
    })
    summary3.render(slide, data)

    print("✅ SummaryComponent test complete")


def run_all_tests():
    """Run all component tests and generate test PowerPoint."""
    print("=" * 60)
    print("ReportForge Component Test Suite")
    print("=" * 60)

    # Create presentation
    prs = Presentation()

    # Prepare test data
    data = create_sample_data()
    trend_data = create_trend_data()

    # Slide 1: Text, Table, Image test
    print("\n--- Slide 1: Text, Table, Image ---")
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    test_text_component(slide1)
    test_table_component(slide1, data)
    test_image_component(slide1)

    # Slide 2: Charts test
    print("\n--- Slide 2: Charts (Column, Pie, Bar) ---")
    test_chart_component_slide(prs, data)

    # Slide 3: Line chart test
    print("\n--- Slide 3: Line Chart ---")
    test_line_chart_slide(prs, trend_data)

    # Slide 4: Summary test
    print("\n--- Slide 4: Summary Component ---")
    test_summary_component_slide(prs, data, trend_data)

    # Save presentation
    output_path = 'output/test_components_output.pptx'
    prs.save(output_path)

    print("\n" + "=" * 60)
    print("✅ ALL TESTS COMPLETE!")
    print(f"Output saved to: {output_path}")
    print("=" * 60)
    print("\nGenerated slides:")
    print("  1. Text, Table, and Image components")
    print("  2. Column, Pie, and Bar charts")
    print("  3. Line chart with multiple series")
    print("  4. Summary component with auto-insights")
    print("\nOpen the file in PowerPoint to view the results!")


if __name__ == '__main__':
    run_all_tests()
