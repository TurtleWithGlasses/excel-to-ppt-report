from pptx import Presentation
from pptx.util import Inches, Pt

def create_table_slide(slide, data):
    """Add a table to the slide with processed data."""
    # Define table position and size
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    rows, cols = data.shape
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Add column headers
    for col_idx, col_name in enumerate(data.columns):
        table.cell(0, col_idx).text = str(col_name)
        cell = table.cell(0, col_idx)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Add data rows
    for row_idx, row in data.iterrows():
        for col_idx, value in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(value)

    print("Table added to the slide.")
