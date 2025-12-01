"""
Main module for ReportForge - Automated PPT Report Generator
Supports both template-based and direct report generation
"""

from data_processing import load_and_process_data
from ppt_generator import create_table_slide
from template_manager import TemplateManager
from pptx import Presentation
import os


def generate_report_from_template(template_id_or_path: str, excel_path: str, output_path: str = None):
    """
    Generate a report using a template configuration
    
    Args:
        template_id_or_path: Template ID or path to template JSON file
        excel_path: Path to Excel data file
        output_path: Optional output path for generated PPT
        
    Returns:
        Path to generated report or None if failed
    """
    # Load template
    template_manager = TemplateManager()
    template = template_manager.load_template(template_id_or_path)
    
    if not template:
        print(f"Failed to load template: {template_id_or_path}")
        return None
    
    print(f"Using template: {template['name']} ({template['client']})")
    
    # Get data mapping from template
    data_mapping = template.get('data_mapping', {})
    sheet_name = data_mapping.get('sheet_name', 'Sheet1')
    
    # Load and process Excel data
    print(f"Loading data from: {excel_path}")
    processed_data = load_and_process_data(excel_path, sheet_name)
    
    if processed_data is None:
        print("No data processed. Exiting...")
        return None
    
    # Load PowerPoint template
    ppt_template_path = template.get('ppt_template_path', '')
    
    if not ppt_template_path or not os.path.exists(ppt_template_path):
        print(f"PPT template not found: {ppt_template_path}")
        return None
    
    print(f"Loading PPT template: {ppt_template_path}")
    presentation = Presentation(ppt_template_path)
    
    # Process slides according to template configuration
    slides_config = template.get('slides', [])
    
    if not slides_config:
        # Fallback to default behavior
        print("No slide configuration found, using default slide")
        slide = presentation.slides[0]
        create_table_slide(slide, processed_data)
    else:
        # Process each configured slide
        for slide_config in slides_config:
            slide_index = slide_config.get('slide_index', 0)
            
            # Ensure slide exists
            if slide_index < len(presentation.slides):
                slide = presentation.slides[slide_index]
                
                # Process components for this slide
                components = slide_config.get('components', [])
                for component in components:
                    component_type = component.get('type')
                    
                    if component_type == 'table':
                        # Create table with specified columns
                        columns = component.get('columns', processed_data.columns.tolist())
                        table_data = processed_data[columns] if columns else processed_data
                        create_table_slide(slide, table_data)
                    
                    # Add more component types (charts, text) in future versions
                    # elif component_type == 'chart':
                    #     create_chart_slide(slide, processed_data, component)
                    # elif component_type == 'text':
                    #     create_text_component(slide, component)
    
    # Determine output path
    if not output_path:
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)
        template_name = template['name'].replace(' ', '_')
        output_path = os.path.join(output_dir, f"{template_name}_Report.pptx")
    
    # Save the updated presentation
    presentation.save(output_path)
    print(f"✓ Report generated successfully: {output_path}")
    
    return output_path


def generate_report_direct(excel_path: str, ppt_template_path: str, sheet_name: str, output_path: str):
    """
    Generate report directly without using a template configuration
    (Legacy method for backward compatibility)
    
    Args:
        excel_path: Path to Excel file
        ppt_template_path: Path to PowerPoint template
        sheet_name: Name of Excel sheet to process
        output_path: Output path for generated report
        
    Returns:
        Path to generated report or None if failed
    """
    print("Generating report (direct mode - no template)")
    
    # Load and process Excel data
    processed_data = load_and_process_data(excel_path, sheet_name)
    if processed_data is None:
        print("No data processed. Exiting...")
        return None

    # Load PowerPoint template
    presentation = Presentation(ppt_template_path)

    # Add processed data to the first slide
    slide = presentation.slides[0]
    create_table_slide(slide, processed_data)

    # Save the updated presentation
    presentation.save(output_path)
    print(f"✓ Report generated: {output_path}")
    
    return output_path


def main():
    """Main entry point - can be used for testing"""
    print("ReportForge - Automated PPT Report Generator")
    print("=" * 50)
    
    # Example 1: Direct generation (legacy method)
    print("\n[Example 1] Direct generation:")
    excel_path = "data/Kasım 2024 BSH ve Rakipleri Medya Analizi Dosyası_.xlsx"
    ppt_template_path = "templates/BSH Kasım Ayı Aylık Medya Yansıma Raporu 24.pptx"
    output_path = "output/Updated_Report.pptx"
    sheet_name = "BSH"
    
    generate_report_direct(excel_path, ppt_template_path, sheet_name, output_path)
    
    # Example 2: Template-based generation
    # Uncomment this when you have templates created
    # print("\n[Example 2] Template-based generation:")
    # generate_report_from_template(
    #     template_id_or_path="templates/configs/your_template.json",
    #     excel_path="data/your_data.xlsx",
    #     output_path="output/Template_Report.pptx"
    # )


if __name__ == "__main__":
    main()
