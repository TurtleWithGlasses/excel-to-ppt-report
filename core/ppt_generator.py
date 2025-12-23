"""
PPTGenerator - Main PowerPoint generation engine

This is the core engine that brings everything together:
- Loads templates
- Loads and maps data
- Creates components
- Generates PowerPoint presentations

This is the main entry point for report generation.
"""

from typing import Dict, Any, List, Optional, Union
from pptx import Presentation
from pptx.util import Inches
import os
from datetime import datetime
import logging

from core.component_factory import ComponentFactory
from core.data_mapper import DataMapper
from core.template_manager import TemplateManager

logger = logging.getLogger(__name__)


class PPTGenerator:
    """
    Main PowerPoint report generation engine.

    Orchestrates the entire report generation process:
    1. Load template
    2. Load and map data
    3. Create presentation
    4. Add slides
    5. Render components
    6. Save PowerPoint file

    Example:
        generator = PPTGenerator()
        generator.load_template('templates/configs/BSH_Template.json')
        generator.load_data('data/BSH_November.xlsx')
        output_path = generator.generate('output/BSH_Report.pptx')
    """

    def __init__(
        self,
        template_dir: str = "templates/configs",
        output_dir: str = "output"
    ):
        """
        Initialize PPTGenerator.

        Args:
            template_dir: Directory containing template JSON files
            output_dir: Directory for generated PowerPoint files
        """
        self.template_manager = TemplateManager(template_dir)
        self.data_mapper = DataMapper()
        self.component_factory = ComponentFactory()

        self.output_dir = output_dir
        self.presentation: Optional[Presentation] = None
        self.custom_variables: Dict[str, Any] = {}

        # Create output directory
        os.makedirs(output_dir, exist_ok=True)

    def load_template(self, template_path: str) -> Dict[str, Any]:
        """
        Load a PowerPoint template.

        Args:
            template_path: Path to template JSON file

        Returns:
            Template dictionary

        Example:
            generator.load_template('templates/configs/BSH_Template.json')
        """
        template = self.template_manager.load_template(template_path)
        return template

    def load_data(self, data_path: str, sheet_name: Union[str, int] = 0) -> None:
        """
        Load data from Excel/CSV file.

        Args:
            data_path: Path to data file
            sheet_name: Sheet name or index for Excel files

        Example:
            generator.load_data('data/BSH_November.xlsx')
        """
        self.data_mapper.load_data(data_path, sheet_name)

    def set_variables(self, variables: Dict[str, Any]) -> None:
        """
        Set custom variables for text substitution.

        Args:
            variables: Dictionary of variables

        Example:
            generator.set_variables({
                'company': 'BSH',
                'month': 'November',
                'report_type': 'Monthly Media Analysis'
            })
        """
        self.custom_variables.update(variables)

    def generate(
        self,
        output_path: Optional[str] = None,
        template_pptx: Optional[str] = None
    ) -> str:
        """
        Generate PowerPoint presentation.

        Args:
            output_path: Path to save generated PowerPoint
                        If None, auto-generates based on template name
            template_pptx: Path to PowerPoint template file (.pptx)
                          If None, creates blank presentation

        Returns:
            Path to generated PowerPoint file

        Example:
            path = generator.generate('output/BSH_November_Report.pptx')
        """
        # Check if template is loaded
        if not self.template_manager.current_template:
            raise ValueError("No template loaded. Call load_template() first.")

        # Create presentation
        if template_pptx and os.path.exists(template_pptx):
            self.presentation = Presentation(template_pptx)
        else:
            self.presentation = Presentation()
            self._set_presentation_size()

        # Get template
        template = self.template_manager.current_template

        # Generate slides
        slides = template.get('slides', [])
        for slide_config in slides:
            self._generate_slide(slide_config)

        # Determine output path
        if not output_path:
            output_path = self._generate_output_path(template)

        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # Save presentation
        self.presentation.save(output_path)

        return output_path

    def _set_presentation_size(self) -> None:
        """Set presentation slide size based on template settings."""
        if not self.presentation:
            return

        template = self.template_manager.current_template
        settings = template.get('settings', {})
        page_size = settings.get('page_size', '16:9')

        # Standard PowerPoint sizes
        if page_size == '16:9':
            self.presentation.slide_width = Inches(10)
            self.presentation.slide_height = Inches(5.625)
        elif page_size == '4:3':
            self.presentation.slide_width = Inches(10)
            self.presentation.slide_height = Inches(7.5)

    def _generate_slide(self, slide_config: Dict[str, Any]) -> None:
        """
        Generate a single slide.

        Args:
            slide_config: Slide configuration from template
        """
        if not self.presentation:
            raise ValueError("Presentation not initialized")

        # Get layout
        layout_name = slide_config.get('layout', 'blank')
        layout_index = self._get_layout_index(layout_name)

        # Add slide
        slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Get components - check both old and new template formats
        components_config = slide_config.get('components', [])

        # Handle new Template Builder format with slide_settings
        if not components_config and 'slide_settings' in slide_config:
            components_config = self._convert_slide_settings_to_components(slide_config)

        # Render each component
        for component_config in components_config:
            self._render_component(slide, component_config)

    def _convert_slide_settings_to_components(self, slide_config: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Convert Template Builder slide_settings format to components array.

        This handles the newer template format where slide configuration is stored
        in slide_settings instead of components array.

        Args:
            slide_config: Slide configuration with slide_settings

        Returns:
            List of component configurations
        """
        components = []
        slide_settings = slide_config.get('slide_settings', {})
        slide_type = slide_config.get('type', '')

        # Handle Table Slide
        if slide_type == 'Table Slide' and 'table' in slide_settings:
            table_settings = slide_settings['table']

            # Add title if present
            title = table_settings.get('title', '')
            if title:
                table_style = table_settings.get('style', {})
                components.append({
                    'type': 'text',
                    'content': title,
                    'position': {'x': 0.5, 'y': 0.3},
                    'size': {'width': 9.0, 'height': 0.6},
                    'style': {
                        'font_name': table_style.get('font_name', 'Calibri'),
                        'font_size': 18,
                        'bold': True,
                        'alignment': 'left'
                    }
                })

            # Get the full style from slide_settings
            slide_table_style = table_settings.get('style', {})

            # Auto-detect computed columns from requested columns
            # These are columns that don't exist in raw data but are computed during aggregation
            computed_column_names = ['Toplam', 'Pozitif', 'Negatif', 'Nötr', 'Kurum',
                                     'YÜKSEK', 'ORTA', 'DÜŞÜK',
                                     'Basın', 'Radyo', 'Televizyon', 'İnternet',
                                     'Ulusal', 'Yerel']
            requested_columns = table_settings.get('columns', [])
            computed_columns = table_settings.get('computed_columns', [])

            # Add any computed columns that are in the requested columns but not already in computed_columns
            for col in requested_columns:
                if col in computed_column_names and col not in computed_columns:
                    computed_columns.append(col)

            # Add table component with ALL settings from the template builder
            component = {
                'type': 'table',
                'position': {'x': 0.5, 'y': 1.0 if title else 0.5},
                'size': {'width': 9.0, 'height': 5.0 if title else 5.5},
                'data_source': {
                    'columns': requested_columns,
                    'sort_by': table_settings.get('sort_by'),
                    'ascending': table_settings.get('ascending', False),
                    'group_by': table_settings.get('group_by'),
                    'aggregations': table_settings.get('aggregations', {}),
                    'computed_columns': computed_columns,
                },
                'style': slide_table_style.copy() if slide_table_style else {}
            }
            components.append(component)

        # Handle Detail Table Slide (shows all rows without grouping/aggregation)
        elif slide_type == 'Detail Table Slide' and 'table' in slide_settings:
            table_settings = slide_settings['table']

            # Add title if present
            title = table_settings.get('title', '')
            if title:
                table_style = table_settings.get('style', {})
                components.append({
                    'type': 'text',
                    'content': title,
                    'position': {'x': 0.5, 'y': 0.3},
                    'size': {'width': 9.0, 'height': 0.6},
                    'style': {
                        'font_name': table_style.get('font_name', 'Calibri'),
                        'font_size': 18,
                        'bold': True,
                        'alignment': 'left'
                    }
                })

            # Get the full style from slide_settings
            slide_table_style = table_settings.get('style', {})

            # Add table component WITHOUT group_by - this shows all rows
            component = {
                'type': 'table',
                'position': {'x': 0.5, 'y': 1.0 if title else 0.5},
                'size': {'width': 9.0, 'height': 5.0 if title else 5.5},
                'data_source': {
                    'columns': table_settings.get('columns', []),
                    'sort_by': table_settings.get('sort_by'),
                    'ascending': table_settings.get('ascending', False),
                    # NO group_by - keep all original rows
                },
                'style': slide_table_style.copy() if slide_table_style else {}
            }
            components.append(component)

        # Handle Chart Slide
        elif slide_type == 'Chart Slide' and 'chart' in slide_settings:
            chart_settings = slide_settings['chart']

            # Add title if present
            title = chart_settings.get('title', '')
            chart_style = chart_settings.get('style', {})
            if title:
                components.append({
                    'type': 'text',
                    'content': title,
                    'position': {'x': 0.5, 'y': 0.3},
                    'size': {'width': 9.0, 'height': 0.6},
                    'style': {
                        'font_name': chart_style.get('font_name', 'Calibri'),
                        'font_size': 18,
                        'bold': True,
                        'alignment': 'left'
                    }
                })

            # Determine if y_column is a computed column
            y_column = chart_settings.get('y_column')
            computed_column_names = ['Toplam', 'Pozitif', 'Negatif', 'Nötr', 'YÜKSEK', 'ORTA', 'DÜŞÜK',
                                     'Basın', 'Radyo', 'Televizyon', 'İnternet', 'Ulusal', 'Yerel']
            computed_columns = []
            if y_column in computed_column_names:
                computed_columns.append(y_column)

            # Get chart type
            chart_type = chart_settings.get('chart_type', 'column')

            # For stacked charts, the y_column may actually be the series column
            # (what to stack by), and we need to count occurrences
            series_column = chart_settings.get('series_column')
            actual_y_column = y_column

            # Check if this is a stacked/grouped chart and y_column is categorical (not numeric)
            if chart_type in ('stacked_bar', 'stacked_column', 'grouped_bar', 'grouped_column'):
                # For stacked/grouped charts, if y_column is a categorical field, use it as series_column
                # and use 'Toplam' (count) as the actual y value
                # Extended list of categorical columns that can be used for grouping/stacking
                categorical_columns = [
                    # Media type columns
                    'Medya Tür', 'Mecra Tipi', 'Yayın Tipi', 'Yayın Türü', 'Mecra',
                    'Medya Kapsam', 'Medya Peryod', 'Medya İçerik', 'Kupür Tipi',
                    # Sentiment/perception columns
                    'Duygu', 'Ton', 'Algı',
                    # Weight/visibility/size columns
                    'Bahis Ağırlığı', 'Görünürlük', 'Boyut',
                    # Location and category columns
                    'Kategori', 'Şehir', 'Medya Şehir',
                    # Other categorical columns
                    'Görsel Malzeme', 'Medya Grup Adı'
                ]
                # Check with whitespace tolerance (strip both y_column and categorical names)
                y_column_stripped = y_column.strip() if y_column else ''
                categorical_columns_stripped = [c.strip() for c in categorical_columns]
                if y_column_stripped in categorical_columns_stripped and not series_column:
                    series_column = y_column
                    actual_y_column = 'Toplam'  # Use count
                    if 'Toplam' not in computed_columns:
                        computed_columns.append('Toplam')

            # Add chart component with ALL settings from the template builder
            component = {
                'type': 'chart',
                'position': {'x': 0.5, 'y': 1.0 if title else 0.5},
                'size': {'width': 9.0, 'height': 5.0 if title else 5.5},
                'data_source': {
                    'x_column': chart_settings.get('x_column'),
                    'y_column': actual_y_column,
                    'series_column': series_column,  # For stacked charts
                    'calculation': chart_settings.get('calculation', 'sum'),
                    'sort_by': chart_settings.get('sort_by'),
                    'ascending': chart_settings.get('ascending', False),
                    'top_n': chart_settings.get('top_n'),
                    'group_by': chart_settings.get('x_column'),  # For chart, group by x column
                    'computed_columns': computed_columns,  # Add computed columns for chart
                },
                'chart_type': chart_type,
                'style': chart_style.copy() if chart_style else {}
            }
            components.append(component)

        # Handle Title Slide
        elif slide_type == 'Title Slide' and 'title' in slide_settings:
            title_settings = slide_settings['title']
            template = self.template_manager.current_template
            if template and 'settings' in template:
                logo_path = template['settings'].get('logo_path')
                embedded_logo_path = template['settings'].get('embedded_logo_path')

                # Add main logo
                if logo_path:
                    logo_pos = title_settings.get('logo_position', {'x': 5.0, 'y': 1.0})
                    logo_size = title_settings.get('logo_size', {'width': 2.0, 'height': 1.5})
                    components.append({
                        'type': 'image',
                        'position': logo_pos,
                        'size': logo_size,
                        'data_source': {'path': logo_path}
                    })

                # Add embedded logo
                if embedded_logo_path:
                    embedded_pos = title_settings.get('embedded_logo_position', {'x': 0.5, 'y': 6.5})
                    embedded_size = title_settings.get('embedded_logo_size', {'width': 1.5, 'height': 0.5})
                    components.append({
                        'type': 'image',
                        'position': embedded_pos,
                        'size': embedded_size,
                        'data_source': {'path': embedded_logo_path}
                    })

                # Add title text
                title_pos = title_settings.get('title_position', {'x': 0.5, 'y': 2.5})
                components.append({
                    'type': 'text',
                    'content': title_settings.get('title', ''),
                    'position': title_pos,
                    'size': {'width': 9.0, 'height': 1.5},
                    'style': {
                        'font_size': 44,
                        'bold': True,
                        'alignment': 'center'
                    }
                })

                # Add subtitle if present
                if title_settings.get('subtitle'):
                    components.append({
                        'type': 'text',
                        'content': title_settings.get('subtitle', ''),
                        'position': {'x': 0.5, 'y': 4.0},
                        'size': {'width': 9.0, 'height': 0.8},
                        'style': {
                            'font_size': 28,
                            'alignment': 'center'
                        }
                    })

                # Add description if present
                if title_settings.get('description'):
                    components.append({
                        'type': 'text',
                        'content': title_settings.get('description', ''),
                        'position': {'x': 0.5, 'y': 4.8},
                        'size': {'width': 9.0, 'height': 1.2},
                        'style': {
                            'font_size': 18,
                            'alignment': 'center'
                        }
                    })

        # Handle Blank Slide - can contain chart, table, or both
        elif slide_type == 'Blank Slide':
            # Check for chart in slide_settings
            if 'chart' in slide_settings:
                chart_settings = slide_settings['chart']
                title = chart_settings.get('title', '')
                chart_style = chart_settings.get('style', {})

                if title:
                    components.append({
                        'type': 'text',
                        'content': title,
                        'position': {'x': 0.5, 'y': 0.3},
                        'size': {'width': 9.0, 'height': 0.6},
                        'style': {
                            'font_name': chart_style.get('font_name', 'Calibri'),
                            'font_size': 18,
                            'bold': True,
                            'alignment': 'left'
                        }
                    })

                # Determine y_column handling for categorical columns
                y_column = chart_settings.get('y_column')
                computed_column_names = ['Toplam', 'Pozitif', 'Negatif', 'Nötr', 'YÜKSEK', 'ORTA', 'DÜŞÜK',
                                         'Basın', 'Radyo', 'Televizyon', 'İnternet', 'Ulusal', 'Yerel']
                computed_columns = []
                if y_column in computed_column_names:
                    computed_columns.append(y_column)

                chart_type = chart_settings.get('chart_type', 'column')
                series_column = chart_settings.get('series_column')
                actual_y_column = y_column

                # Check if this is a stacked/grouped chart and y_column is categorical
                if chart_type in ('stacked_bar', 'stacked_column', 'grouped_bar', 'grouped_column'):
                    categorical_columns = [
                        'Medya Tür', 'Mecra Tipi', 'Yayın Tipi', 'Yayın Türü', 'Mecra',
                        'Medya Kapsam', 'Medya Peryod', 'Medya İçerik', 'Kupür Tipi',
                        'Duygu', 'Ton', 'Algı',
                        'Bahis Ağırlığı', 'Görünürlük', 'Boyut',
                        'Kategori', 'Şehir', 'Medya Şehir',
                        'Görsel Malzeme', 'Medya Grup Adı'
                    ]
                    y_column_stripped = y_column.strip() if y_column else ''
                    categorical_columns_stripped = [c.strip() for c in categorical_columns]
                    if y_column_stripped in categorical_columns_stripped and not series_column:
                        series_column = y_column
                        actual_y_column = 'Toplam'
                        if 'Toplam' not in computed_columns:
                            computed_columns.append('Toplam')

                component = {
                    'type': 'chart',
                    'position': {'x': 0.5, 'y': 1.0 if title else 0.5},
                    'size': {'width': 9.0, 'height': 5.0 if title else 5.5},
                    'data_source': {
                        'x_column': chart_settings.get('x_column'),
                        'y_column': actual_y_column,
                        'series_column': series_column,
                        'calculation': chart_settings.get('calculation', 'sum'),
                        'sort_by': chart_settings.get('sort_by'),
                        'ascending': chart_settings.get('ascending', False),
                        'top_n': chart_settings.get('top_n'),
                        'group_by': chart_settings.get('x_column'),
                        'computed_columns': computed_columns,
                    },
                    'chart_type': chart_type,
                    'style': chart_style.copy() if chart_style else {}
                }
                components.append(component)

            # Check for table in slide_settings
            if 'table' in slide_settings:
                table_settings = slide_settings['table']
                title = table_settings.get('title', '')
                table_style = table_settings.get('style', {})

                # Only add title if no chart title was added
                if title and 'chart' not in slide_settings:
                    components.append({
                        'type': 'text',
                        'content': title,
                        'position': {'x': 0.5, 'y': 0.3},
                        'size': {'width': 9.0, 'height': 0.6},
                        'style': {
                            'font_name': table_style.get('font_name', 'Calibri'),
                            'font_size': 18,
                            'bold': True,
                            'alignment': 'left'
                        }
                    })

                # Auto-detect computed columns for Blank Slide tables
                computed_column_names = ['Toplam', 'Pozitif', 'Negatif', 'Nötr', 'Kurum',
                                         'YÜKSEK', 'ORTA', 'DÜŞÜK',
                                         'Basın', 'Radyo', 'Televizyon', 'İnternet',
                                         'Ulusal', 'Yerel']
                requested_columns = table_settings.get('columns', [])
                computed_columns = table_settings.get('computed_columns', [])
                for col in requested_columns:
                    if col in computed_column_names and col not in computed_columns:
                        computed_columns.append(col)

                component = {
                    'type': 'table',
                    'position': {'x': 0.5, 'y': 1.0 if title else 0.5},
                    'size': {'width': 9.0, 'height': 5.0 if title else 5.5},
                    'data_source': {
                        'columns': requested_columns,
                        'sort_by': table_settings.get('sort_by'),
                        'ascending': table_settings.get('ascending', False),
                        'group_by': table_settings.get('group_by'),
                        'aggregations': table_settings.get('aggregations', {}),
                        'computed_columns': computed_columns,
                    },
                    'style': table_style.copy() if table_style else {}
                }
                components.append(component)

        return components

    def _get_layout_index(self, layout_name: str) -> int:
        """
        Get slide layout index by name.

        Args:
            layout_name: Layout name (title, blank, etc.)

        Returns:
            Layout index
        """
        layout_map = {
            'title': 0,
            'title_content': 1,
            'section': 2,
            'two_content': 3,
            'comparison': 4,
            'blank': 5,
            'content': 6
        }

        return layout_map.get(layout_name.lower(), 5)  # Default to blank

    def _render_component(self, slide, component_config: Dict[str, Any]) -> None:
        """
        Render a component on a slide.

        Args:
            slide: PowerPoint slide object
            component_config: Component configuration
        """
        try:
            # Create component
            component = self.component_factory.create_component(component_config)

            # Get data for component
            component_data = self.data_mapper.get_data_for_component(
                component_config,
                self.custom_variables
            )

            # Render component
            component.render(slide, component_data)

        except Exception as e:
            print(f"Warning: Failed to render component: {str(e)}")
            # Continue with other components

    def _generate_output_path(self, template: Dict[str, Any]) -> str:
        """
        Generate output path based on template name and timestamp.

        Args:
            template: Template dictionary

        Returns:
            Output file path
        """
        metadata = template.get('metadata', {})
        template_name = metadata.get('name', 'Report')

        # Clean filename
        clean_name = template_name.replace(' ', '_').replace('/', '_')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        filename = f"{clean_name}_{timestamp}.pptx"
        return os.path.join(self.output_dir, filename)

    def generate_from_config(
        self,
        template_path: str,
        data_path: str,
        output_path: Optional[str] = None,
        variables: Optional[Dict[str, Any]] = None,
        sheet_name: Union[str, int] = 0
    ) -> str:
        """
        Generate PowerPoint from template and data (convenience method).

        Args:
            template_path: Path to template JSON file
            data_path: Path to data file (Excel/CSV)
            output_path: Path to save generated PowerPoint
            variables: Custom variables for text substitution
            sheet_name: Sheet name or index for Excel files

        Returns:
            Path to generated PowerPoint file

        Example:
            generator = PPTGenerator()
            output = generator.generate_from_config(
                template_path='templates/configs/BSH_Template.json',
                data_path='data/BSH_November.xlsx',
                variables={'company': 'BSH', 'month': 'November'}
            )
        """
        # Load template
        self.load_template(template_path)

        # Load data
        self.load_data(data_path, sheet_name)

        # Set variables
        if variables:
            self.set_variables(variables)

        # Generate
        return self.generate(output_path)

    def get_template_info(self) -> Dict[str, Any]:
        """
        Get information about loaded template.

        Returns:
            Template information dictionary
        """
        return self.template_manager.get_template_info()

    def get_data_info(self) -> Dict[str, Any]:
        """
        Get information about loaded data.

        Returns:
            Data information dictionary
        """
        return self.data_mapper.get_info()

    def validate_template(self) -> bool:
        """
        Validate current template.

        Returns:
            True if valid

        Raises:
            ValueError: If template is invalid
        """
        return self.template_manager.validate_current_template()

    def preview_variables(self) -> Dict[str, Any]:
        """
        Preview available variables for text substitution.

        Returns:
            Dictionary of all available variables
        """
        return self.data_mapper.create_variable_dict(self.custom_variables)

    def reset(self) -> None:
        """Reset generator to initial state."""
        self.presentation = None
        self.custom_variables = {}
        self.data_mapper.reset()


class BatchPPTGenerator:
    """
    Batch PowerPoint generation for multiple reports.

    Example:
        batch_gen = BatchPPTGenerator()
        batch_gen.add_job(
            template='templates/configs/BSH_Template.json',
            data='data/BSH_November.xlsx',
            variables={'company': 'BSH'}
        )
        batch_gen.add_job(
            template='templates/configs/Sanofi_Template.json',
            data='data/Sanofi_November.xlsx',
            variables={'company': 'Sanofi'}
        )
        results = batch_gen.generate_all()
    """

    def __init__(self, output_dir: str = "output"):
        """
        Initialize BatchPPTGenerator.

        Args:
            output_dir: Directory for generated PowerPoint files
        """
        self.output_dir = output_dir
        self.jobs: List[Dict[str, Any]] = []
        self.results: List[Dict[str, Any]] = []

    def add_job(
        self,
        template: str,
        data: str,
        output: Optional[str] = None,
        variables: Optional[Dict[str, Any]] = None,
        sheet_name: Union[str, int] = 0
    ) -> None:
        """
        Add a generation job to the batch.

        Args:
            template: Path to template JSON file
            data: Path to data file
            output: Output path (optional)
            variables: Custom variables
            sheet_name: Sheet name or index
        """
        self.jobs.append({
            'template': template,
            'data': data,
            'output': output,
            'variables': variables or {},
            'sheet_name': sheet_name
        })

    def generate_all(self) -> List[Dict[str, Any]]:
        """
        Generate all jobs in the batch.

        Returns:
            List of results with status and output paths
        """
        self.results = []

        for idx, job in enumerate(self.jobs):
            print(f"Processing job {idx + 1}/{len(self.jobs)}...")

            try:
                generator = PPTGenerator(output_dir=self.output_dir)

                output_path = generator.generate_from_config(
                    template_path=job['template'],
                    data_path=job['data'],
                    output_path=job['output'],
                    variables=job['variables'],
                    sheet_name=job['sheet_name']
                )

                self.results.append({
                    'job_index': idx,
                    'status': 'success',
                    'output_path': output_path,
                    'template': job['template'],
                    'data': job['data']
                })

                print(f"✅ Job {idx + 1} completed: {output_path}")

            except Exception as e:
                self.results.append({
                    'job_index': idx,
                    'status': 'failed',
                    'error': str(e),
                    'template': job['template'],
                    'data': job['data']
                })

                print(f"❌ Job {idx + 1} failed: {str(e)}")

        return self.results

    def get_summary(self) -> Dict[str, Any]:
        """
        Get summary of batch generation results.

        Returns:
            Summary dictionary
        """
        if not self.results:
            return {'status': 'No jobs completed'}

        successful = sum(1 for r in self.results if r['status'] == 'success')
        failed = sum(1 for r in self.results if r['status'] == 'failed')

        return {
            'total_jobs': len(self.jobs),
            'successful': successful,
            'failed': failed,
            'success_rate': f"{(successful / len(self.jobs) * 100):.1f}%",
            'results': self.results
        }


# Example usage
if __name__ == '__main__':
    print("=" * 60)
    print("PPTGenerator Test")
    print("=" * 60)

    print("\nThis is a placeholder test.")
    print("Full integration test requires:")
    print("  1. Template JSON file")
    print("  2. Sample Excel data file")
    print("  3. Component library")
    print("\nRun test_core_engine.py for complete testing.")

    print("\n" + "=" * 60)
    print("✅ PPTGenerator Module Loaded!")
    print("=" * 60)
